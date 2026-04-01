#!/usr/bin/env python3
"""Report-like styled PPTX generator with cover/body distinction.

Key behavior:
- Cover slide uses hero layout (dark background)
- Body slides use structured layout (section marker + underline + cards)
- Metadata anchors on every slide:
  - top-right: date_label
  - bottom-left: author_display_name
  - bottom-right: presentation_title
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

COL = {
    "dgistnavy": RGBColor(14, 32, 80),
    "dgistblue": RGBColor(26, 60, 140),
    "dgistmid": RGBColor(60, 100, 180),
    "dgistlight": RGBColor(220, 230, 250),
    "bglight": RGBColor(248, 249, 252),
    "bgpurple": RGBColor(248, 240, 255),
    "bgorange": RGBColor(255, 247, 235),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "line": RGBColor(216, 225, 243),
}
FONT = "Apple SD Gothic Neo"


def add_meta(slide, date_label, author, pres_title, dark=False):
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear()
    p = tr.paragraphs[0]
    p.text = date_label
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 239, 255) if dark else COL["muted"]

    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.0), Inches(0.2)).text_frame
    bl.clear()
    a = bl.paragraphs[0]
    a.text = author
    a.font.name = FONT
    a.font.size = Pt(11)
    a.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]

    br = slide.shapes.add_textbox(Inches(7.0), Inches(7.04), Inches(5.6), Inches(0.2)).text_frame
    br.clear()
    t = br.paragraphs[0]
    t.text = pres_title
    t.alignment = PP_ALIGN.RIGHT
    t.font.name = FONT
    t.font.size = Pt(11)
    t.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]


def add_cover(prs, title, subtitle, date_label, author, pres_title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COL["dgistnavy"]

    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.3))
    top.fill.solid(); top.fill.fore_color.rgb = COL["dgistblue"]
    top.line.fill.background()

    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(11.3), Inches(3.8))
    hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(34, 56, 116)
    hero.fill.transparency = 0.18
    hero.line.color.rgb = RGBColor(160, 182, 236)

    tf = hero.text_frame
    tf.clear(); tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.bold = True; p.font.size = Pt(48); p.font.color.rgb = RGBColor(255, 255, 255)
    q = tf.add_paragraph()
    q.text = subtitle
    q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = COL["dgistlight"]

    add_meta(s, date_label, author, pres_title, dark=True)


def add_body_frame(slide, sec_num, title, subtitle, date_label, author, pres_title):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL["bglight"]

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.33))
    strip.fill.solid(); strip.fill.fore_color.rgb = COL["dgistblue"]
    strip.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.58), Inches(0.62), Inches(0.42))
    marker.fill.solid(); marker.fill.fore_color.rgb = COL["dgistblue"]
    marker.line.fill.background()
    mtf = marker.text_frame; mtf.clear()
    mp = mtf.paragraphs[0]
    mp.text = str(sec_num)
    mp.alignment = PP_ALIGN.CENTER
    mp.font.name = FONT; mp.font.size = Pt(15); mp.font.bold = True; mp.font.color.rgb = RGBColor(255, 255, 255)

    t = slide.shapes.add_textbox(Inches(1.38), Inches(0.56), Inches(10.9), Inches(0.5)).text_frame
    t.clear(); tp = t.paragraphs[0]
    tp.text = title
    tp.font.name = FONT; tp.font.size = Pt(33); tp.font.bold = True; tp.font.color.rgb = COL["dgistnavy"]

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.08), Inches(10.95), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = COL["dgistblue"]
    line.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(1.42), Inches(1.16), Inches(10.7), Inches(0.34)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(16); sp.font.color.rgb = COL["muted"]

    add_meta(slide, date_label, author, pres_title, dark=False)


def add_card(slide, x, y, w, h, title=None, fill="bglight"):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = COL[fill]
    sh.line.color.rgb = COL["line"]
    sh.line.width = Pt(1.0)
    sh.adjustments[0] = 0.05
    if title:
        tt = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.32)).text_frame
        tt.clear(); p = tt.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["dgistblue"]
    return sh


def fill_bullets(shape, bullets, size=20, top=0.58):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(0.22); tf.margin_top = Inches(top)
    for i, b in enumerate(bullets[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = COL["text"]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--title', required=True)
    ap.add_argument('--subtitle', default='')
    ap.add_argument('--author-display-name', default='[author_display_name]')
    ap.add_argument('--presentation-title', default='[presentation_title]')
    ap.add_argument('--date-label', default='[date_label]')
    ap.add_argument('--out', required=True)
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs, args.title, args.subtitle, args.date_label, args.author_display_name, args.presentation_title)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_body_frame(s, 1, 'Overview', '', args.date_label, args.author_display_name, args.presentation_title)
    c = add_card(s, 0.7, 1.9, 12.1, 4.8, 'Summary', 'bglight')
    fill_bullets(c, ['Point 1', 'Point 2', 'Point 3'])

    prs.save(args.out)
    print(args.out)


if __name__ == '__main__':
    main()
