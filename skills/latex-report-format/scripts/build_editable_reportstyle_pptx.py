#!/usr/bin/env python3
"""Build a fully editable PPTX using report-style design tokens (no image embedding)."""

from datetime import date
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

C = {
    "navy": RGBColor(14, 32, 80),
    "blue": RGBColor(26, 60, 140),
    "mid": RGBColor(60, 100, 180),
    "light": RGBColor(220, 230, 250),
    "bg": RGBColor(248, 249, 252),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(216, 225, 243),
    "white": RGBColor(255, 255, 255),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "bgpurple": RGBColor(248, 240, 255),
    "bgorange": RGBColor(255, 247, 235),
    "bggreen": RGBColor(235, 248, 240),
}
FONT = "Apple SD Gothic Neo"


def add_meta(slide, date_label, author, pres_title, dark=False):
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear(); p = tr.paragraphs[0]
    p.text = date_label; p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255) if dark else C["muted"]

    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.2), Inches(0.2)).text_frame
    bl.clear(); q = bl.paragraphs[0]
    q.text = author
    q.font.name = FONT; q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]

    br = slide.shapes.add_textbox(Inches(6.6), Inches(7.04), Inches(6.0), Inches(0.2)).text_frame
    br.clear(); r = br.paragraphs[0]
    r.text = pres_title; r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]


def card(slide, x,y,w,h, fill, title=None, title_color=None):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = C["line"]; s.line.width = Pt(1.0)
    s.adjustments[0] = 0.05
    if title:
        tf = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.12), Inches(w-0.4), Inches(0.32)).text_frame
        tf.clear(); p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True
        p.font.color.rgb = title_color or C["blue"]
    return s


def bullets(shape, lines, size=19, top=0.58):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18); tf.margin_top = Inches(top)
    for i, t in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {t}"
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = C["text"]
        p.space_after = Pt(8)


def body_header(slide, sec, title, subtitle):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = C["blue"]; bar.line.fill.background()

    m = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(0.56), Inches(0.62), Inches(0.42))
    m.fill.solid(); m.fill.fore_color.rgb = C["blue"]; m.line.fill.background()
    mt = m.text_frame; mt.clear(); mp = mt.paragraphs[0]
    mp.text = str(sec); mp.alignment = PP_ALIGN.CENTER
    mp.font.name = FONT; mp.font.size = Pt(15); mp.font.bold = True; mp.font.color.rgb = RGBColor(255,255,255)

    tt = slide.shapes.add_textbox(Inches(1.38), Inches(0.54), Inches(10.95), Inches(0.5)).text_frame
    tt.clear(); p = tt.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = C["navy"]

    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.06), Inches(10.95), Inches(0.03))
    ul.fill.solid(); ul.fill.fore_color.rgb = C["blue"]; ul.line.fill.background()

    st = slide.shapes.add_textbox(Inches(1.42), Inches(1.14), Inches(10.8), Inches(0.34)).text_frame
    st.clear(); sp = st.paragraphs[0]
    sp.text = subtitle
    sp.font.name = FONT; sp.font.size = Pt(15.5); sp.font.color.rgb = C["muted"]


def build(out, author, pres_title, date_label):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C["navy"]
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.30))
    top.fill.solid(); top.fill.fore_color.rgb = C["blue"]; top.line.fill.background()

    org = s.shapes.add_textbox(Inches(4.1), Inches(0.95), Inches(5.2), Inches(0.35)).text_frame
    org.clear(); p = org.paragraphs[0]
    p.text = "DGIST / RESEARCH SEMINAR"; p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = C["light"]

    main = card(s, 2.0, 1.45, 9.35, 2.3, RGBColor(255,255,255), None)
    main.fill.transparency = 0.92
    tf = main.text_frame; tf.clear(); tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.38)
    p1 = tf.paragraphs[0]; p1.text = "Attention Is All You Need"; p1.alignment = PP_ALIGN.CENTER
    p1.font.name = FONT; p1.font.size = Pt(44); p1.font.bold = True; p1.font.color.rgb = RGBColor(255,255,255)
    p2 = tf.add_paragraph(); p2.text = "Transformer Architecture Brief"; p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT; p2.font.size = Pt(20); p2.font.color.rgb = C["light"]

    t3 = s.shapes.add_textbox(Inches(3.15), Inches(4.1), Inches(7.1), Inches(0.5)).text_frame
    t3.clear(); q = t3.paragraphs[0]
    q.text = "2026년 03월 발표 자료"; q.alignment = PP_ALIGN.CENTER
    q.font.name = FONT; q.font.size = Pt(25); q.font.bold = True; q.font.color.rgb = C["light"]

    info = s.shapes.add_textbox(Inches(3.25), Inches(4.85), Inches(6.9), Inches(1.35)).text_frame
    info.clear()
    for i,row in enumerate([f"발표자: {author}", f"발표명: {pres_title}", "소속: DGIST"]):
        rp = info.paragraphs[0] if i == 0 else info.add_paragraph()
        rp.text = row; rp.alignment = PP_ALIGN.CENTER
        rp.font.name = FONT; rp.font.size = Pt(16); rp.font.color.rgb = C["light"]
    add_meta(s, date_label, author, pres_title, dark=True)

    # slide 2 overview
    s = prs.slides.add_slide(prs.slide_layouts[6]); body_header(s,1,"연구 개요","한 문장 요약과 핵심 메시지")
    add_meta(s, date_label, author, pres_title, dark=False)
    key = card(s,0.72,1.92,12.0,2.05,C["bg"] ,"핵심 요약")
    bullets(key,["RNN/CNN 없이 Self-Attention 중심 구조를 제안","학습 병렬화를 크게 향상해 대규모 학습 기반을 마련","현대 LLM 아키텍처의 출발점 역할"],21,0.55)
    l = card(s,0.72,4.15,5.85,2.55,C["bgorange"],"문제의식",C["blue"])
    bullets(l,["장거리 의존성 학습 어려움","순차 처리의 병렬화 한계","학습 비용 증가"],18)
    r = card(s,6.87,4.15,5.85,2.55,C["bgpurple"],"핵심 전환",C["blue"])
    bullets(r,["Attention-only 구조 채택","직접적 토큰 관계 계산","확장 가능한 학습 패턴 제공"],18)

    # slide 3 mechanism
    s = prs.slides.add_slide(prs.slide_layouts[6]); body_header(s,2,"핵심 연산","Scaled Dot-Product Attention")
    add_meta(s, date_label, author, pres_title, dark=False)
    panel = card(s,0.72,1.92,12.0,4.78,C["card"],"Attention formula")
    eq = s.shapes.add_textbox(Inches(1.0),Inches(2.55),Inches(11.4),Inches(0.8)).text_frame
    eq.clear(); p=eq.paragraphs[0]
    p.text="Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name=FONT; p.font.size=Pt(30); p.font.bold=True; p.font.color.rgb=C["blue"]
    box = s.shapes.add_textbox(Inches(1.0),Inches(3.55),Inches(11.4),Inches(2.3)).text_frame
    box.clear()
    for i,t in enumerate(["QK^T로 토큰 간 관련도 계산","sqrt(d_k) 스케일링으로 softmax 안정화","V 가중합으로 문맥 표현 생성"]):
        q = box.paragraphs[0] if i==0 else box.add_paragraph()
        q.text=f"• {t}"; q.font.name=FONT; q.font.size=Pt(21); q.font.color.rgb=C["text"]

    # slide 4 results table
    s = prs.slides.add_slide(prs.slide_layouts[6]); body_header(s,3,"결과와 영향","보고서 스타일 표 구조")
    add_meta(s, date_label, author, pres_title, dark=False)
    card(s,0.72,1.92,12.0,4.78,C["card"],None)
    table = s.shapes.add_table(4,2,Inches(1.0),Inches(2.28),Inches(11.4),Inches(3.72)).table
    table.columns[0].width=Inches(2.8); table.columns[1].width=Inches(8.6)
    table.cell(0,0).text='Category'; table.cell(0,1).text='Impact'
    for cidx in range(2):
        cell=table.cell(0,cidx); cell.fill.solid(); cell.fill.fore_color.rgb=C["blue"]
        pp=cell.text_frame.paragraphs[0]; pp.font.name=FONT; pp.font.size=Pt(14); pp.font.bold=True; pp.font.color.rgb=C["white"]
    rows=[("Modeling","Recurrence 중심에서 Attention 중심으로 전환"),("Efficiency","학습 병렬화 및 확장성 대폭 향상"),("Legacy","BERT/GPT 등 현대 LLM 구조의 기반 확립")]
    for i,(k,v) in enumerate(rows, start=1):
        table.cell(i,0).text=k; table.cell(i,1).text=v
        for cidx in range(2):
            pp=table.cell(i,cidx).text_frame.paragraphs[0]
            pp.font.name=FONT; pp.font.size=Pt(13.5); pp.font.color.rgb=C["text"]

    # slide 5 limitations + takeaways (detail refinement)
    s = prs.slides.add_slide(prs.slide_layouts[6]); body_header(s,4,"한계와 요약","발표 마무리용 핵심 포인트")
    add_meta(s, date_label, author, pres_title, dark=False)
    l = card(s,0.72,1.92,5.85,4.78,C["bgorange"],"한계",C["blue"])
    bullets(l,["Self-Attention의 계산 비용은 O(n^2)","긴 문맥에서 메모리/지연 비용 증가","배포 환경에서 효율 최적화 필요"],18)
    r = card(s,6.87,1.92,5.85,4.78,C["bgpurple"],"요약",C["blue"])
    bullets(r,["Transformer 이해는 LLM 이해의 출발점","성능과 효율을 함께 최적화해야 함","Attention 설계 선택이 모델 역량을 좌우"],18)

    prs.save(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--out', required=True)
    ap.add_argument('--author-display-name', default='Jaeyeong CHOI')
    ap.add_argument('--presentation-title', default='Attention Is All You Need Presentation')
    ap.add_argument('--date-label', default=str(date.today()))
    args = ap.parse_args()
    build(args.out, args.author_display_name, args.presentation_title, args.date_label)
    print(args.out)


if __name__ == '__main__':
    main()
