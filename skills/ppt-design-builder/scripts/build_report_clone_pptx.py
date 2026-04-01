#!/usr/bin/env python3
"""Create PPTX by placing each report PDF page as a centered image per slide.

Use this mode when user asks to keep report design *as-is* and simply shrink into PPT.
"""

import argparse
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image


def convert_pdf_to_pngs(pdf: Path, out_dir: Path, dpi: int = 240):
    out_dir.mkdir(parents=True, exist_ok=True)
    prefix = out_dir / 'page'
    subprocess.check_call(['pdftoppm', '-png', '-r', str(dpi), str(pdf), str(prefix)])
    return sorted(out_dir.glob('page-*.png'), key=lambda p: int(p.stem.split('-')[-1]))


def build_ppt_from_images(images, out_pptx: Path, full_bleed: bool = True):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 252)

        if full_bleed:
            slide.shapes.add_picture(str(img), 0, 0, width=sw, height=sh)
            continue

        with Image.open(img) as im:
            w, h = im.size

        margin_x = Inches(0.55)
        margin_y = Inches(0.30)
        max_w = sw - 2 * margin_x
        max_h = sh - 2 * margin_y

        ratio = w / h
        box_ratio = max_w / max_h
        if ratio > box_ratio:
            draw_w = max_w
            draw_h = max_w / ratio
        else:
            draw_h = max_h
            draw_w = max_h * ratio

        left = (sw - draw_w) / 2
        top = (sh - draw_h) / 2
        slide.shapes.add_picture(str(img), left, top, width=int(draw_w), height=int(draw_h))

    prs.save(out_pptx)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--pdf', required=True, help='Input report PDF path')
    ap.add_argument('--out', required=True, help='Output PPTX path')
    ap.add_argument('--tmp-dir', default='./tmp_report_pages', help='Temporary PNG output directory')
    ap.add_argument('--dpi', type=int, default=240)
    ap.add_argument('--fit-inside', action='store_true',
                    help='Keep page inside margins (disables full-bleed)')
    args = ap.parse_args()

    pdf = Path(args.pdf)
    out = Path(args.out)
    tmp = Path(args.tmp_dir)

    images = convert_pdf_to_pngs(pdf, tmp, dpi=args.dpi)
    build_ppt_from_images(images, out, full_bleed=(not args.fit_inside))
    print(out)


if __name__ == '__main__':
    main()
