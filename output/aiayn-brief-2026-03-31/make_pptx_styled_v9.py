from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

# Report-aligned palette
C = {
    "navy": RGBColor(14, 32, 80),
    "blue": RGBColor(26, 60, 140),
    "mid": RGBColor(60, 100, 180),
    "light": RGBColor(220, 230, 250),
    "bg": RGBColor(248, 249, 252),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(216, 225, 243),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "soft1": RGBColor(248, 240, 255),
    "soft2": RGBColor(255, 247, 235),
    "soft3": RGBColor(235, 248, 240),
}
FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_meta(slide, dark=False):
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear(); p = tr.paragraphs[0]
    p.text = DATE_LABEL; p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255) if dark else C["muted"]

    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.0), Inches(0.2)).text_frame
    bl.clear(); q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT; q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]

    br = slide.shapes.add_textbox(Inches(6.7), Inches(7.04), Inches(5.9), Inches(0.2)).text_frame
    br.clear(); r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE; r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]


def cover(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C["navy"]

    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    top.fill.solid(); top.fill.fore_color.rgb = C["blue"]; top.line.fill.background()

    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.25), Inches(11.45), Inches(4.55))
    hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(255, 255, 255); hero.fill.transparency = 0.88
    hero.line.color.rgb = RGBColor(175, 196, 242)

    tf = hero.text_frame
    tf.clear(); tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.55)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.size = Pt(50); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)

    q = tf.add_paragraph()
    q.text = subtitle
    q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = C["light"]

    chips = tf.add_paragraph()
    chips.text = "Self-Attention  ·  Parallel Training  ·  LLM Foundation"
    chips.font.name = FONT; chips.font.size = Pt(16); chips.font.color.rgb = RGBColor(207, 221, 255)

    add_meta(s, dark=True)


def body_header(slide, idx, title, subtitle=""):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = C["blue"]; bar.line.fill.background()

    num = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(0.58), Inches(0.62), Inches(0.42))
    num.fill.solid(); num.fill.fore_color.rgb = C["blue"]; num.line.fill.background()
    nt = num.text_frame; nt.clear(); n = nt.paragraphs[0]
    n.text = str(idx); n.alignment = PP_ALIGN.CENTER
    n.font.name = FONT; n.font.size = Pt(15); n.font.bold = True; n.font.color.rgb = RGBColor(255, 255, 255)

    tt = slide.shapes.add_textbox(Inches(1.38), Inches(0.56), Inches(10.9), Inches(0.5)).text_frame
    tt.clear(); t = tt.paragraphs[0]
    t.text = title; t.font.name = FONT; t.font.size = Pt(33); t.font.bold = True; t.font.color.rgb = C["navy"]

    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.08), Inches(10.95), Inches(0.03))
    ul.fill.solid(); ul.fill.fore_color.rgb = C["blue"]; ul.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(1.42), Inches(1.16), Inches(10.7), Inches(0.34)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(16); sp.font.color.rgb = C["muted"]

    add_meta(slide, dark=False)


def card(slide, x, y, w, h, title=None, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill or C["card"]
    shape.line.color.rgb = C["line"]; shape.line.width = Pt(1.05)
    shape.adjustments[0] = 0.05
    if title:
        ht = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.12), Inches(w-0.4), Inches(0.32)).text_frame
        ht.clear(); p = ht.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = C["blue"]
    return shape


def bullets(shape, lines, size=20, top=0.58):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18); tf.margin_top = Inches(top)
    for i, line in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = C["text"]
        p.space_after = Pt(8)


def slide_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 1, "Why this paper mattered", "from recurrence bottlenecks to attention-first modeling")
    l = card(s, 0.72, 1.92, 6.0, 4.8, "Before", C["card"])
    r = card(s, 6.92, 1.92, 5.68, 4.8, "After", C["soft2"])
    bullets(l, [
        "Long-range dependency learning was difficult",
        "Sequential computation limited parallelism",
        "Scaling up caused high training cost",
    ], 20)
    bullets(r, [
        "Attention directly models token relations",
        "Architecture is naturally parallelizable",
        "Foundation for modern large language models",
    ], 20)


def slide_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 2, "Architecture summary", "encoder/decoder + self-attention blocks")
    panel = card(s, 0.72, 1.92, 11.88, 4.8, "Pipeline", C["card"])
    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    fills = [C["soft1"], C["soft3"], C["soft1"], C["soft3"]]
    for i, lbl in enumerate(labels):
        x = 1.16 + i * 2.8
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.34), Inches(2.3), Inches(1.24))
        b.fill.solid(); b.fill.fore_color.rgb = fills[i]
        b.line.color.rgb = C["line"]
        tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = lbl; p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = C["navy"]
    for i in range(3):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5 + i*2.8), Inches(3.72), Inches(0.42), Inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb = C["mid"]
        a.line.fill.background()


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 3, "Core mechanism", "scaled dot-product attention")
    panel = card(s, 0.72, 1.92, 11.88, 4.8, "Attention formula", C["card"])
    eq = s.shapes.add_textbox(Inches(1.02), Inches(2.56), Inches(11.3), Inches(0.78)).text_frame
    eq.clear(); p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = C["blue"]
    ex = s.shapes.add_textbox(Inches(1.02), Inches(3.54), Inches(11.3), Inches(2.4)).text_frame
    ex.clear()
    for i, t in enumerate([
        "QK^T computes token-to-token relevance",
        "sqrt(d_k) scaling improves optimization stability",
        "Weighted V aggregation builds contextual representation",
    ]):
        q = ex.paragraphs[0] if i == 0 else ex.add_paragraph()
        q.text = f"• {t}"; q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = C["text"]


def slide_results():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 4, "Results and impact", "report-style table block")

    # outer panel
    panel = card(s, 0.72, 1.92, 11.88, 4.8, None, C["card"])

    # table
    table = s.shapes.add_table(4, 2, Inches(1.02), Inches(2.30), Inches(11.30), Inches(3.65)).table
    table.columns[0].width = Inches(2.75)
    table.columns[1].width = Inches(8.55)

    table.cell(0,0).text = "Category"
    table.cell(0,1).text = "Impact"
    for c in range(2):
        cell = table.cell(0,c)
        cell.fill.solid(); cell.fill.fore_color.rgb = C["blue"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

    rows = [
        ("Modeling", "Shift from recurrence-centric to attention-centric sequence modeling"),
        ("Efficiency", "Substantially improved parallel training characteristics"),
        ("Legacy", "Established backbone used by BERT/GPT and modern LLM families"),
    ]
    for i,(k,v) in enumerate(rows, start=1):
        table.cell(i,0).text = k
        table.cell(i,1).text = v
        for c in range(2):
            p = table.cell(i,c).text_frame.paragraphs[0]
            p.font.name = FONT; p.font.size = Pt(13.5); p.font.color.rgb = C["text"]


def slide_limits():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 5, "Limitations and follow-ups", "what remains challenging")
    l = card(s, 0.72, 1.92, 6.0, 4.8, "Known limits", C["card"])
    r = card(s, 6.92, 1.92, 5.68, 4.8, "Follow-up lines", C["soft3"])
    bullets(l, [
        "Self-attention cost grows quadratically with sequence length",
        "Long-context settings increase memory and latency",
        "Deployment requires careful efficiency optimization",
    ], 20)
    bullets(r, [
        "Sparse / Linear Attention variants",
        "FlashAttention-style kernel optimizations",
        "Long-context specialized architecture design",
    ], 20)


def slide_takeaways():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 6, "Takeaways", "three points to remember")
    panel = card(s, 0.72, 1.92, 11.88, 4.8, "Final summary", C["soft1"])
    tf = panel.text_frame
    tf.clear(); tf.margin_left = Inches(0.32); tf.margin_top = Inches(0.72)
    points = [
        "Transformer literacy is core LLM literacy.",
        "Model quality and efficiency should be optimized together.",
        "Attention design choices strongly shape capability.",
    ]
    for i,txt in enumerate(points):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = f"{i+1}. {txt}"
        p.font.name = FONT; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C["navy"]
        p.space_after = Pt(12)


cover("Attention Is All You Need", "report-template visual language adapted for slides")
slide_problem()
slide_architecture()
slide_equation()
slide_results()
slide_limits()
slide_takeaways()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v9.pptx"
prs.save(out)
print(out)
