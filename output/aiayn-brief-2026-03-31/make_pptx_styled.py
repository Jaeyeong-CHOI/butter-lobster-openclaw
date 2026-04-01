from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = int(13.333 * 914400)  # 16:9
prs.slide_height = int(7.5 * 914400)

COLORS = {
    "bg": RGBColor(246, 248, 254),
    "title": RGBColor(18, 41, 99),
    "accent": RGBColor(33, 78, 166),
    "body": RGBColor(38, 40, 52),
    "sub": RGBColor(95, 103, 130),
}

FONT_TITLE = "Aptos Display"
FONT_BODY = "Aptos"

slides = [
    ("Attention Is All You Need", ["Design-aware concise briefing", "Transformer의 출발점"]),
    ("Why it was needed", ["RNN Seq2Seq had bottlenecks", "Long-range dependency was hard", "Training parallelism was limited"]),
    ("Core architecture", ["Encoder-Decoder Transformer", "Multi-Head Self-Attention", "FFN + Residual + LayerNorm"]),
    ("Attention formula", ["Attention(Q,K,V)=softmax(QK^T/sqrt(d_k))V", "Scaling improves optimization stability", "Token relationships are computed directly"]),
    ("Multi-head + Position", ["Different heads learn different relations", "Positional encoding injects order", "No recurrence needed"]),
    ("Impact", ["Strong translation performance", "Better parallel training efficiency", "Foundation for BERT/GPT-style LLMs"]),
    ("Limitations", ["Self-attention cost grows as O(n^2)", "Long context is expensive", "Follow-ups: sparse/linear/flash attention"]),
    ("Takeaways", ["Transformer is core LLM literacy", "Quality and efficiency must be co-optimized", "Attention design shapes capability"]),
]

for idx, (title, bullets) in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]

    # Top accent bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, int(0.32 * 914400))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent"]
    bar.line.fill.background()

    # Title
    tbox = slide.shapes.add_textbox(int(0.7 * 914400), int(0.45 * 914400), int(11.8 * 914400), int(1.1 * 914400))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(42 if idx == 0 else 36)
    p.font.color.rgb = COLORS["title"]

    # Bullet area card
    card = slide.shapes.add_shape(1, int(0.75 * 914400), int(1.7 * 914400), int(11.8 * 914400), int(4.9 * 914400))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(222, 228, 246)

    btf = card.text_frame
    btf.clear()
    for j, b in enumerate(bullets):
        bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        bp.text = b
        bp.level = 0
        bp.font.name = FONT_BODY
        bp.font.size = Pt(24 if idx == 0 else 23)
        bp.font.color.rgb = COLORS["body"]

    # Footer note
    fbox = slide.shapes.add_textbox(int(0.8 * 914400), int(6.95 * 914400), int(12.0 * 914400), int(0.35 * 914400))
    ft = fbox.text_frame
    ft.clear()
    fp = ft.paragraphs[0]
    fp.text = "OpenClaw · Attention Is All You Need Brief"
    fp.font.name = FONT_BODY
    fp.font.size = Pt(12)
    fp.font.color.rgb = COLORS["sub"]

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled.pptx"
prs.save(out)
print(out)
