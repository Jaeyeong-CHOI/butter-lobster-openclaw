from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

COL = {
    "dgistnavy": RGBColor(14, 32, 80),
    "dgistblue": RGBColor(26, 60, 140),
    "dgistmid": RGBColor(60, 100, 180),
    "dgistlight": RGBColor(220, 230, 250),
    "accentgreen": RGBColor(39, 130, 67),
    "accentorange": RGBColor(210, 100, 20),
    "accentpurple": RGBColor(100, 50, 150),
    "bglight": RGBColor(248, 249, 252),
    "bgpurple": RGBColor(248, 240, 255),
    "bgorange": RGBColor(255, 247, 235),
    "bggray": RGBColor(245, 245, 247),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "line": RGBColor(216, 225, 243),
}
FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_meta(slide, dark=False):
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear(); p = tr.paragraphs[0]
    p.text = DATE_LABEL; p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255) if dark else COL["muted"]

    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.0), Inches(0.2)).text_frame
    bl.clear(); q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT; q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]

    br = slide.shapes.add_textbox(Inches(7.0), Inches(7.04), Inches(5.6), Inches(0.2)).text_frame
    br.clear(); r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE; r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]


def cover_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = COL["dgistnavy"]

    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.3))
    top.fill.solid(); top.fill.fore_color.rgb = COL["dgistblue"]; top.line.fill.background()

    badge = s.shapes.add_textbox(Inches(0.85), Inches(1.05), Inches(5.0), Inches(0.3)).text_frame
    badge.clear(); bp = badge.paragraphs[0]
    bp.text = "PROJECT / PAPER BRIEF"
    bp.font.name = FONT; bp.font.size = Pt(12); bp.font.bold = True; bp.font.color.rgb = COL["dgistlight"]

    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.4), Inches(11.65), Inches(3.95))
    hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(30, 52, 108); hero.fill.transparency = 0.16
    hero.line.color.rgb = RGBColor(161, 184, 236)

    tf = hero.text_frame
    tf.clear(); tf.margin_left = Inches(0.48); tf.margin_top = Inches(0.58)
    p = tf.paragraphs[0]
    p.text = "Attention Is All You Need"
    p.font.name = FONT; p.font.size = Pt(50); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
    q = tf.add_paragraph()
    q.text = "Transformer architecture brief in report-style design"
    q.font.name = FONT; q.font.size = Pt(20); q.font.color.rgb = COL["dgistlight"]

    k = tf.add_paragraph()
    k.text = "Self-Attention · Parallel Training · LLM Foundation"
    k.font.name = FONT; k.font.size = Pt(17); k.font.color.rgb = RGBColor(202, 219, 255)

    add_meta(s, dark=True)


def body_frame(slide, num, title, subtitle=""):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = COL["bglight"]

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.33))
    strip.fill.solid(); strip.fill.fore_color.rgb = COL["dgistblue"]; strip.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.58), Inches(0.62), Inches(0.42))
    marker.fill.solid(); marker.fill.fore_color.rgb = COL["dgistblue"]; marker.line.fill.background()
    mt = marker.text_frame; mt.clear(); mp = mt.paragraphs[0]
    mp.text = str(num); mp.alignment = PP_ALIGN.CENTER
    mp.font.name = FONT; mp.font.size = Pt(15); mp.font.bold = True; mp.font.color.rgb = RGBColor(255, 255, 255)

    t = slide.shapes.add_textbox(Inches(1.38), Inches(0.56), Inches(10.9), Inches(0.5)).text_frame
    t.clear(); tp = t.paragraphs[0]
    tp.text = title; tp.font.name = FONT; tp.font.size = Pt(33); tp.font.bold = True; tp.font.color.rgb = COL["dgistnavy"]

    uline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.08), Inches(10.95), Inches(0.03))
    uline.fill.solid(); uline.fill.fore_color.rgb = COL["dgistblue"]; uline.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(1.42), Inches(1.16), Inches(10.7), Inches(0.34)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle; sp.font.name = FONT; sp.font.size = Pt(16); sp.font.color.rgb = COL["muted"]

    add_meta(slide, dark=False)


def card(slide, x, y, w, h, title=None, fill="bggray", title_color="dgistblue"):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = COL[fill]
    sh.line.color.rgb = COL["line"]; sh.line.width = Pt(1.0)
    sh.adjustments[0] = 0.05
    if title:
        tt = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.12), Inches(w-0.4), Inches(0.32)).text_frame
        tt.clear(); p = tt.paragraphs[0]
        p.text = title; p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL[title_color]
    return sh


def put_bullets(shape, lines, size=20, top=0.58):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(0.22); tf.margin_top = Inches(top)
    for i, line in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"; p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = COL["text"]
        p.space_after = Pt(8)


def slide_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 2, "Why this paper mattered", "from recurrence bottlenecks to attention-first modeling")
    left = card(s, 0.7, 1.9, 6.1, 4.8, "Before", "bggray")
    right = card(s, 6.95, 1.9, 5.85, 4.8, "After", "bgorange", "accentorange")
    put_bullets(left, [
        "Long-range dependency learning was difficult",
        "Sequential computation limited parallelism",
        "Scaling caused high training cost",
    ], 20)
    put_bullets(right, [
        "Attention directly models token relations",
        "Architecture is naturally parallelizable",
        "Enabled modern large-scale language models",
    ], 20)


def slide_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 3, "Architecture summary", "encoder/decoder + self-attention blocks")
    panel = card(s, 0.7, 1.9, 12.1, 4.75, "Pipeline", "bggray")
    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    for i, lbl in enumerate(labels):
        x = 1.2 + i * 2.86
        block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.35), Inches(2.35), Inches(1.25))
        block.fill.solid(); block.fill.fore_color.rgb = COL["bgpurple"] if i % 2 == 0 else COL["dgistlight"]
        block.line.color.rgb = COL["line"]
        tf = block.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = lbl; p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["dgistnavy"]
    for i in range(3):
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.58 + i * 2.86), Inches(3.73), Inches(0.45), Inches(0.38))
        arr.fill.solid(); arr.fill.fore_color.rgb = COL["dgistmid"]
        arr.line.fill.background()


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 4, "Core mechanism", "scaled dot-product attention")
    panel = card(s, 0.7, 1.9, 12.1, 4.8, "Attention formula", "bggray")
    eq = s.shapes.add_textbox(Inches(1.0), Inches(2.58), Inches(11.7), Inches(0.78)).text_frame
    eq.clear(); p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = COL["dgistblue"]
    ex = s.shapes.add_textbox(Inches(1.0), Inches(3.55), Inches(11.7), Inches(2.35)).text_frame
    ex.clear()
    for i, t in enumerate([
        "QK^T computes token-to-token relevance",
        "sqrt(d_k) scaling improves optimization stability",
        "Weighted V aggregation builds contextual representation",
    ]):
        q = ex.paragraphs[0] if i == 0 else ex.add_paragraph()
        q.text = f"• {t}"; q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = COL["text"]


def slide_impact_table():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 5, "Results and impact", "report-like table block")
    # table block with colored header row
    outer = card(s, 0.7, 1.9, 12.1, 4.8, None, "bggray")
    table = s.shapes.add_table(4, 2, Inches(1.0), Inches(2.35), Inches(11.5), Inches(3.6)).table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(8.7)

    # header
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Impact"
    for c in range(2):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = COL["dgistblue"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)

    rows = [
        ("Modeling", "Shift from recurrence-centric to attention-centric sequence modeling"),
        ("Efficiency", "Substantially improved parallel training characteristics"),
        ("Legacy", "Established backbone used by BERT/GPT and modern LLM families"),
    ]
    for i, (k, v) in enumerate(rows, start=1):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
        for c in range(2):
            p = table.cell(i, c).text_frame.paragraphs[0]
            p.font.name = FONT; p.font.size = Pt(13.5); p.font.color.rgb = COL["text"]


def slide_takeaway():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 6, "Takeaways", "three points to remember")
    panel = card(s, 0.7, 1.95, 12.1, 4.8, "Final summary", "bgpurple", "accentpurple")
    tf = panel.text_frame
    tf.clear(); tf.margin_left = Inches(0.32); tf.margin_top = Inches(0.72)
    for i, txt in enumerate([
        "Transformer literacy is core LLM literacy.",
        "Model quality and computational efficiency must be evaluated together.",
        "Attention design choices strongly shape model capability.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {txt}"
        p.font.name = FONT; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = COL["dgistnavy"]
        p.space_after = Pt(12)


cover_slide()
slide_problem()
slide_architecture()
slide_equation()
slide_impact_table()
slide_takeaway()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v8.pptx"
prs.save(out)
print(out)
