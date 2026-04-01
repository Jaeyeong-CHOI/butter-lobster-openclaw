from pathlib import Path
import subprocess
from pptx import Presentation
from pptx.util import Inches

pdf = Path('/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_aiayn_reportstyle_16x9.pdf')
outdir = Path('/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/tmp_v15_pages')
outppt = Path('/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_report-identical_v15.pptx')
outdir.mkdir(parents=True, exist_ok=True)

subprocess.check_call(['pdftoppm', '-png', '-r', '300', str(pdf), str(outdir / 'page')])
imgs = sorted(outdir.glob('page-*.png'), key=lambda p: int(p.stem.split('-')[-1]))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for img in imgs:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(outppt)
print(outppt)
