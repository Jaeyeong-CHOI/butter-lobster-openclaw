from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

# refined report-like palette
C = {
    "navy": RGBColor(14, 32, 80),
    "blue": RGBColor(26, 60, 140),
    "mid": RGBColor(60, 100, 180),
    "bg": RGBColor(248, 249, 252),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(216, 225, 243),
    "text": RGBColor(34, 38, 50),
    "muted": RGBColor(97, 109, 140),
    "soft_blue": RGBColor(238, 244, 255),
    "soft_purple": RGBColor(247, 240, 255),
    "soft_green": RGBColor(236, 248, 241),
}

FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_meta(slide, dark=False):
    # top-right date
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.06), Inches(3.0), Inches(0.2)).text_frame
    tr.clear()
    p = tr.paragraphs[0]
    p.text = DATE_LABEL
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 238, 255) if dark else C["muted"]

    # bottom-left author
    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.4), Inches(0.2)).text_frame
    bl.clear()
    q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT
    q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]

    # bottom-right presentation title
    br = slide.shapes.add_textbox(Inches(6.4), Inches(7.04), Inches(6.2), Inches(0.2)).text_frame
    br.clear()
    r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE
    r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]


def cover_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navy"]

    # top ribbon
    ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    ribbon.fill.solid(); ribbon.fill.fore_color.rgb = C["blue"]
    ribbon.line.fill.background()

    # central hero panel
    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.28), Inches(11.45), Inches(4.55))
    hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.fill.transparency = 0.9
    hero.line.color.rgb = RGBColor(162, 184, 236)
    hero.line.width = Pt(1.2)

    tf = hero.text_frame
    tf.clear()
    tf.margin_left = Inches(0.55)
    tf.margin_top = Inches(0.62)

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.bold = True
    p.font.size = Pt(50)
    p.font.color.rgb = RGBColor(255, 255, 255)

    q = tf.add_paragraph()
    q.text = subtitle
    q.font.name = FONT
    q.font.size = Pt(21)
    q.font.color.rgb = RGBColor(212, 224, 255)

    k = tf.add_paragraph()
    k.text = "Self-Attention · Parallel Training · LLM Foundation"
    k.font.name = FONT
    k.font.size = Pt(16)
    k.font.color.rgb = RGBColor(187, 205, 247)

    add_meta(s, dark=True)


def body_frame(slide, sec_num, title, subtitle=""):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.30))
    bar.fill.solid(); bar.fill.fore_color.rgb = C["blue"]
    bar.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.56), Inches(0.62), Inches(0.42))
    marker.fill.solid(); marker.fill.fore_color.rgb = C["blue"]
    marker.line.fill.background()

    mt = marker.text_frame
    mt.clear()
    m = mt.paragraphs[0]
    m.text = str(sec_num)
    m.alignment = PP_ALIGN.CENTER
    m.font.name = FONT; m.font.size = Pt(15); m.font.bold = True
    m.font.color.rgb = C["white"]

    tt = slide.shapes.add_textbox(Inches(1.38), Inches(0.54), Inches(10.95), Inches(0.5)).text_frame
    tt.clear(); t = tt.paragraphs[0]
    t.text = title
    t.font.name = FONT; t.font.size = Pt(33); t.font.bold = True
    t.font.color.rgb = C["navy"]

    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.06), Inches(10.95), Inches(0.03))
    ul.fill.solid(); ul.fill.fore_color.rgb = C["blue"]
    ul.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(1.42), Inches(1.14), Inches(10.8), Inches(0.34)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(16)
        sp.font.color.rgb = C["muted"]

    add_meta(slide, dark=False)


def card(slide, x, y, w, h, title=None, fill=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill or C["white"]
    c.line.color.rgb = C["line"]
    c.line.width = Pt(1.05)
    c.adjustments[0] = 0.05

    if title:
        ht = slide.shapes.add_textbox(Inches(x+0.22), Inches(y+0.12), Inches(w-0.45), Inches(0.32)).text_frame
        ht.clear(); p = ht.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True
        p.font.color.rgb = C["blue"]
    return c


def bullets(shape, lines, size=20, top=0.58):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(0.1)

    for i, line in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = C["text"]
        p.space_after = Pt(8)


def slide_problem_context():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 1, "Problem and context", "why sequence modeling needed a new architecture")

    before = card(s, 0.72, 1.90, 6.0, 4.80, "Before", C["white"])
    after = card(s, 6.92, 1.90, 5.68, 4.80, "After", C["soft_blue"])

    bullets(before, [
        "RNN recurrence made long-range dependencies hard to learn",
        "Sequential processing limited hardware parallelism",
        "Scaling model/data significantly increased training cost",
    ], size=19)

    bullets(after, [
        "Self-attention directly models token-to-token interactions",
        "Parallel-friendly architecture improved throughput",
        "Enabled practical large-scale pretraining pipelines",
    ], size=19)


def slide_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 2, "Architecture overview", "encoder/decoder with repeated attention blocks")

    panel = card(s, 0.72, 1.92, 11.88, 4.78, "Transformer pipeline", C["white"])

    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    fills = [C["soft_purple"], C["soft_green"], C["soft_purple"], C["soft_green"]]

    for i, lbl in enumerate(labels):
        x = 1.16 + i * 2.8
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.32), Inches(2.30), Inches(1.22))
        b.fill.solid(); b.fill.fore_color.rgb = fills[i]
        b.line.color.rgb = C["line"]
        tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = lbl
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True
        p.font.color.rgb = C["navy"]

    for i in range(3):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.50 + i*2.8), Inches(3.71), Inches(0.42), Inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb = C["mid"]
        a.line.fill.background()


def slide_mechanism():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 3, "Core mechanism", "scaled dot-product attention")

    panel = card(s, 0.72, 1.90, 11.88, 4.80, "Attention formula", C["white"])

    eq = s.shapes.add_textbox(Inches(1.02), Inches(2.56), Inches(11.30), Inches(0.78)).text_frame
    eq.clear(); p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = C["blue"]

    ex = s.shapes.add_textbox(Inches(1.02), Inches(3.52), Inches(11.30), Inches(2.35)).text_frame
    ex.clear()
    for i, t in enumerate([
        "QK^T computes pairwise relevance between tokens",
        "sqrt(d_k) scaling prevents unstable softmax saturation",
        "Weighted V aggregation forms contextual token representations",
    ]):
        q = ex.paragraphs[0] if i == 0 else ex.add_paragraph()
        q.text = f"• {t}"
        q.font.name = FONT
        q.font.size = Pt(21)
        q.font.color.rgb = C["text"]


def slide_results_impact():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 4, "Results and impact", "what changed after AIAYN")

    panel = card(s, 0.72, 1.90, 11.88, 4.80, None, C["white"])

    table = s.shapes.add_table(4, 2, Inches(1.02), Inches(2.28), Inches(11.30), Inches(3.68)).table
    table.columns[0].width = Inches(2.78)
    table.columns[1].width = Inches(8.52)

    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Impact"
    for c in range(2):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = C["blue"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C["white"]

    rows = [
        ("Modeling", "Shifted from recurrence-centric to attention-centric sequence modeling"),
        ("Efficiency", "Substantially improved training parallelism and scalability"),
        ("Legacy", "Established core backbone used by BERT, GPT, and modern LLM families"),
    ]

    for i, (k, v) in enumerate(rows, start=1):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
        for c in range(2):
            p = table.cell(i, c).text_frame.paragraphs[0]
            p.font.name = FONT
            p.font.size = Pt(13.5)
            p.font.color.rgb = C["text"]


def slide_limits_and_takeaways():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_frame(s, 5, "Limitations and takeaways", "practical constraints and next design lessons")

    left = card(s, 0.72, 1.90, 6.0, 4.80, "Limitations", C["white"])
    right = card(s, 6.92, 1.90, 5.68, 4.80, "Takeaways", C["soft_purple"])

    bullets(left, [
        "Self-attention complexity still grows quadratically with sequence length",
        "Long-context settings increase memory footprint and latency",
        "Production deployment needs aggressive efficiency optimization",
    ], size=19)

    bullets(right, [
        "Transformer literacy is essential for LLM work",
        "Quality and efficiency must be optimized together",
        "Attention design choices strongly shape downstream capability",
    ], size=19)


cover_slide("Attention Is All You Need", "refined, report-aligned presentation design")
slide_problem_context()
slide_architecture()
slide_mechanism()
slide_results_impact()
slide_limits_and_takeaways()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v10.pptx"
prs.save(out)
print(out)
