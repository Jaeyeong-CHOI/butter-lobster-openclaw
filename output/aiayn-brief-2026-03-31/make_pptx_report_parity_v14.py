from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

# Colors copied from report template tokens
C = {
    "dgistnavy": RGBColor(14, 32, 80),
    "dgistblue": RGBColor(26, 60, 140),
    "dgistmid": RGBColor(60, 100, 180),
    "dgistlight": RGBColor(220, 230, 250),
    "accentgreen": RGBColor(39, 130, 67),
    "accentorange": RGBColor(210, 100, 20),
    "accentpurple": RGBColor(100, 50, 150),
    "bglight": RGBColor(248, 249, 252),
    "card": RGBColor(255, 255, 255),
    "bggreen": RGBColor(235, 248, 240),
    "bgorange": RGBColor(255, 247, 235),
    "bgpurple": RGBColor(248, 240, 255),
    "tabhead": RGBColor(26, 60, 140),
    "line": RGBColor(216, 225, 243),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "white": RGBColor(255, 255, 255),
}

FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_meta(slide, dark=False):
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear(); p = tr.paragraphs[0]
    p.text = DATE_LABEL
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255) if dark else C["muted"]

    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(5.0), Inches(0.2)).text_frame
    bl.clear(); q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT; q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]

    br = slide.shapes.add_textbox(Inches(6.6), Inches(7.04), Inches(6.0), Inches(0.2)).text_frame
    br.clear(); r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE
    r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT; r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else C["muted"]


def cover_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C["dgistnavy"]

    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.30))
    top.fill.solid(); top.fill.fore_color.rgb = C["dgistblue"]
    top.line.fill.background()

    org = s.shapes.add_textbox(Inches(4.2), Inches(0.95), Inches(5.0), Inches(0.35)).text_frame
    org.clear(); p = org.paragraphs[0]
    p.text = "ORGANIZATION / PROGRAM"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = C["dgistlight"]

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.45), Inches(9.35), Inches(2.3))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(255, 255, 255); box.fill.transparency = 0.92
    box.line.color.rgb = RGBColor(170, 190, 235)

    tf = box.text_frame
    tf.clear(); tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.35)
    t1 = tf.paragraphs[0]
    t1.text = "Attention Is All You Need"
    t1.alignment = PP_ALIGN.CENTER
    t1.font.name = FONT; t1.font.size = Pt(44); t1.font.bold = True; t1.font.color.rgb = C["white"]
    t2 = tf.add_paragraph()
    t2.text = "Transformer Architecture Brief"
    t2.alignment = PP_ALIGN.CENTER
    t2.font.name = FONT; t2.font.size = Pt(20); t2.font.color.rgb = C["dgistlight"]

    t3 = s.shapes.add_textbox(Inches(3.15), Inches(4.1), Inches(7.1), Inches(0.5)).text_frame
    t3.clear(); q = t3.paragraphs[0]
    q.text = "2026년 03월 연구 발표"
    q.alignment = PP_ALIGN.CENTER
    q.font.name = FONT; q.font.size = Pt(25); q.font.bold = True; q.font.color.rgb = C["dgistlight"]

    info = s.shapes.add_textbox(Inches(3.3), Inches(4.85), Inches(6.8), Inches(1.35)).text_frame
    info.clear()
    rows = [
        "작성: Jaeyeong CHOI",
        "소속: DGIST",
        "주제: Attention Is All You Need",
    ]
    for i, row in enumerate(rows):
        rp = info.paragraphs[0] if i == 0 else info.add_paragraph()
        rp.text = row
        rp.alignment = PP_ALIGN.CENTER
        rp.font.name = FONT; rp.font.size = Pt(16); rp.font.color.rgb = C["dgistlight"]

    add_meta(s, dark=True)


def body_header(slide, sec, title, subtitle=""):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bglight"]

    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C["dgistblue"]
    hdr.line.fill.background()

    hl = slide.shapes.add_textbox(Inches(0.35), Inches(0.06), Inches(4.0), Inches(0.2)).text_frame
    hl.clear(); h = hl.paragraphs[0]
    h.text = "PROJECT | ORG"
    h.font.name = FONT; h.font.size = Pt(10.5); h.font.bold = True; h.font.color.rgb = RGBColor(235,240,255)

    # section marker style
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(0.56), Inches(0.62), Inches(0.42))
    marker.fill.solid(); marker.fill.fore_color.rgb = C["dgistblue"]
    marker.line.fill.background()
    mt = marker.text_frame; mt.clear(); mp = mt.paragraphs[0]
    mp.text = str(sec)
    mp.alignment = PP_ALIGN.CENTER
    mp.font.name = FONT; mp.font.size = Pt(15); mp.font.bold = True; mp.font.color.rgb = C["white"]

    title_tf = slide.shapes.add_textbox(Inches(1.38), Inches(0.54), Inches(10.95), Inches(0.5)).text_frame
    title_tf.clear(); tp = title_tf.paragraphs[0]
    tp.text = title
    tp.font.name = FONT; tp.font.size = Pt(32); tp.font.bold = True; tp.font.color.rgb = C["dgistnavy"]

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.06), Inches(10.95), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = C["dgistblue"]
    line.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(1.42), Inches(1.14), Inches(10.8), Inches(0.34)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(15.5); sp.font.color.rgb = C["muted"]

    add_meta(slide, dark=False)


def card(slide, x, y, w, h, title=None, fill=None, title_color=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill or C["card"]
    c.line.color.rgb = C["line"]
    c.line.width = Pt(1.05)
    c.adjustments[0] = 0.05
    if title:
        t = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.12), Inches(w-0.4), Inches(0.32)).text_frame
        t.clear(); p = t.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = title_color or C["dgistblue"]
    return c


def bullets(shape, lines, size=19, top=0.58):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18); tf.margin_top = Inches(top)
    for i, line in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = C["text"]
        p.space_after = Pt(8)


def slide_overview():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 1, "연구 개요", "한 문장 요약과 핵심 메시지")

    key = card(s, 0.72, 1.92, 12.0, 2.05, "핵심 요약", C["bglight"])
    bullets(key, [
        "RNN/CNN 없이 Self-Attention 중심 구조를 제안",
        "학습 병렬화를 크게 향상해 대규모 학습 기반을 마련",
        "현대 LLM 아키텍처의 출발점 역할",
    ], size=21, top=0.55)

    row1 = card(s, 0.72, 4.15, 5.85, 2.55, "문제의식", C["bgorange"], C["accentorange"])
    bullets(row1, [
        "장거리 의존성 학습 어려움",
        "순차 처리의 병렬화 한계",
        "학습 비용 증가",
    ], size=18)

    row2 = card(s, 6.87, 4.15, 5.85, 2.55, "핵심 전환", C["bgpurple"], C["accentpurple"])
    bullets(row2, [
        "Attention-only 구조 채택",
        "직접적 토큰 관계 계산",
        "확장 가능한 학습 패턴 제공",
    ], size=18)


def slide_arch():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 2, "핵심 방법", "Transformer 인코더-디코더 구성")

    panel = card(s, 0.72, 1.92, 12.0, 4.78, "Pipeline", C["card"])
    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    fills = [C["bgpurple"], C["bggreen"], C["bgpurple"], C["bggreen"]]
    for i, lbl in enumerate(labels):
        x = 1.15 + i * 2.80
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.30), Inches(2.30), Inches(1.22))
        b.fill.solid(); b.fill.fore_color.rgb = fills[i]
        b.line.color.rgb = C["line"]
        tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = lbl; p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT; p.font.size = Pt(14.5); p.font.bold = True; p.font.color.rgb = C["dgistnavy"]

    for i in range(3):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.49 + i*2.80), Inches(3.70), Inches(0.42), Inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb = C["dgistmid"]
        a.line.fill.background()

    note = s.shapes.add_textbox(Inches(1.02), Inches(5.0), Inches(11.4), Inches(1.5)).text_frame
    note.clear(); p = note.paragraphs[0]
    p.text = "Residual + LayerNorm + Position-wise FFN 조합으로 안정적 학습 및 표현력 확보"
    p.font.name = FONT; p.font.size = Pt(18); p.font.color.rgb = C["text"]


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 3, "핵심 연산", "Scaled Dot-Product Attention")

    panel = card(s, 0.72, 1.92, 12.0, 4.78, "Attention formula", C["card"])
    eq = s.shapes.add_textbox(Inches(1.00), Inches(2.55), Inches(11.4), Inches(0.80)).text_frame
    eq.clear(); p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = C["dgistblue"]

    bullets_box = s.shapes.add_textbox(Inches(1.00), Inches(3.55), Inches(11.4), Inches(2.25)).text_frame
    bullets_box.clear()
    arr = [
        "QK^T로 토큰 간 관련도 계산",
        "sqrt(d_k) 스케일링으로 softmax 안정화",
        "V 가중합으로 문맥 표현 생성",
    ]
    for i, t in enumerate(arr):
        q = bullets_box.paragraphs[0] if i == 0 else bullets_box.add_paragraph()
        q.text = f"• {t}"
        q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = C["text"]


def slide_results():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 4, "결과와 영향", "보고서 스타일 표 구조")

    card(s, 0.72, 1.92, 12.0, 4.78, None, C["card"])
    table = s.shapes.add_table(4, 2, Inches(1.0), Inches(2.28), Inches(11.4), Inches(3.72)).table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(8.6)

    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Impact"
    for c in range(2):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = C["tabhead"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = C["white"]

    rows = [
        ("Modeling", "Recurrence 중심에서 Attention 중심으로 전환"),
        ("Efficiency", "학습 병렬화 및 확장성 대폭 향상"),
        ("Legacy", "BERT/GPT 등 현대 LLM 구조의 기반 확립"),
    ]
    for i, (k, v) in enumerate(rows, start=1):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
        for c in range(2):
            p = table.cell(i, c).text_frame.paragraphs[0]
            p.font.name = FONT; p.font.size = Pt(13.5); p.font.color.rgb = C["text"]


def slide_takeaway():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 5, "한계와 요약", "실무 관점에서 기억할 핵심")

    l = card(s, 0.72, 1.92, 6.0, 4.78, "한계", C["bgorange"], C["accentorange"])
    r = card(s, 6.92, 1.92, 5.8, 4.78, "요약", C["bgpurple"], C["accentpurple"])

    bullets(l, [
        "Self-Attention 비용은 시퀀스 길이에 따라 O(n^2)",
        "긴 문맥에서 메모리/지연 비용 증가",
        "배포 단계에서 효율 최적화 필수",
    ], 18)

    bullets(r, [
        "Transformer 이해는 LLM 이해의 핵심",
        "성능과 효율을 함께 최적화해야 함",
        "Attention 설계 선택이 모델 역량을 좌우",
    ], 18)


# build
cover_slide()
slide_overview()
slide_arch()
slide_equation()
slide_results()
slide_takeaway()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_report-parity_v14.pptx"
prs.save(out)
print(out)
