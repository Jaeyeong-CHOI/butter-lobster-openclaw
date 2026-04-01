from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

img_dir = Path('/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/report_pages')
out = Path('/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_report-clone_v11.pptx')

images = sorted(img_dir.glob('page-*.png'), key=lambda p: int(p.stem.split('-')[-1]))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_w = prs.slide_width
slide_h = prs.slide_height

for img_path in images:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # subtle report-like background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor(248, 249, 252)

    with Image.open(img_path) as im:
        w, h = im.size

    # fit image inside slide with margins, preserving aspect
    margin_x = Inches(0.55)
    margin_y = Inches(0.30)
    max_w = slide_w - 2 * margin_x
    max_h = slide_h - 2 * margin_y

    img_ratio = w / h
    box_ratio = max_w / max_h

    if img_ratio > box_ratio:
        draw_w = max_w
        draw_h = max_w / img_ratio
    else:
        draw_h = max_h
        draw_w = max_h * img_ratio

    left = (slide_w - draw_w) / 2
    top = (slide_h - draw_h) / 2

    slide.shapes.add_picture(str(img_path), left, top, width=int(draw_w), height=int(draw_h))

prs.save(out)
print(out)
