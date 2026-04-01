from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# metadata
AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

# report-template palette (mirrored)
COL = {
    "dgistnavy": RGBColor(14, 32, 80),
    "dgistblue": RGBColor(26, 60, 140),
    "dgistmid": RGBColor(60, 100, 180),
    "dgistlight": RGBColor(220, 230, 250),
    "accentgreen": RGBColor(39, 130, 67),
    "accentorange": RGBColor(210, 100, 20),
    "accentpurple": RGBColor(100, 50, 150),
    "bglight": RGBColor(248, 249, 252),
    "bgpurple": RGBColor(248, 240, 255),
    "bgorange": RGBColor(255, 247, 235),
    "bggray": RGBColor(245, 245, 247),
    "tabhead": RGBColor(26, 60, 140),
    "text": RGBColor(35, 38, 50),
    "muted": RGBColor(98, 109, 138),
    "line": RGBColor(216, 225, 243),
}

FONT_MAIN = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def base(slide, header_left="PROJECT | ORG"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COL["bglight"]

    # top header strip like report
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.33))
    strip.fill.solid(); strip.fill.fore_color.rgb = COL["dgistblue"]
    strip.line.fill.background()

    # header left
    hl = slide.shapes.add_textbox(Inches(0.35), Inches(0.05), Inches(4.2), Inches(0.2)).text_frame
    hl.clear()
    p = hl.paragraphs[0]
    p.text = header_left
    p.font.name = FONT_MAIN
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255)

    # top-right date (requested)
    tr = slide.shapes.add_textbox(Inches(10.0), Inches(0.05), Inches(3.0), Inches(0.2)).text_frame
    tr.clear()
    q = tr.paragraphs[0]
    q.text = DATE_LABEL
    q.alignment = PP_ALIGN.RIGHT
    q.font.name = FONT_MAIN
    q.font.size = Pt(10.5)
    q.font.bold = True
    q.font.color.rgb = RGBColor(235, 240, 255)

    # bottom-left author (requested)
    bl = slide.shapes.add_textbox(Inches(0.45), Inches(7.04), Inches(4.8), Inches(0.2)).text_frame
    bl.clear()
    a = bl.paragraphs[0]
    a.text = AUTHOR
    a.font.name = FONT_MAIN
    a.font.size = Pt(11)
    a.font.color.rgb = COL["muted"]

    # bottom-right presentation title (requested)
    br = slide.shapes.add_textbox(Inches(7.05), Inches(7.04), Inches(5.8), Inches(0.2)).text_frame
    br.clear()
    t = br.paragraphs[0]
    t.text = PRESENTATION_TITLE
    t.alignment = PP_ALIGN.RIGHT
    t.font.name = FONT_MAIN
    t.font.size = Pt(11)
    t.font.color.rgb = COL["muted"]


def section_title(slide, sec_num, title, subtitle=""):
    # numbered label box like report section style
    num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.58), Inches(0.62), Inches(0.42))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = COL["dgistblue"]
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(sec_num)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_MAIN; p.font.bold = True; p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(255, 255, 255)

    title_tf = slide.shapes.add_textbox(Inches(1.38), Inches(0.56), Inches(10.8), Inches(0.52)).text_frame
    title_tf.clear()
    t = title_tf.paragraphs[0]
    t.text = title
    t.font.name = FONT_MAIN
    t.font.bold = True
    t.font.size = Pt(33)
    t.font.color.rgb = COL["dgistnavy"]

    # blue underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(1.08), Inches(10.95), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = COL["dgistblue"]
    line.line.fill.background()

    if subtitle:
        stf = slide.shapes.add_textbox(Inches(1.42), Inches(1.16), Inches(10.7), Inches(0.34)).text_frame
        stf.clear()
        s = stf.paragraphs[0]
        s.text = subtitle
        s.font.name = FONT_MAIN
        s.font.size = Pt(16)
        s.font.color.rgb = COL["muted"]


def card(slide, x, y, w, h, fill="bglight", title=None, title_color="dgistblue"):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = COL[fill]
    sh.line.color.rgb = COL["line"]
    sh.line.width = Pt(1.1)
    sh.adjustments[0] = 0.05

    if title:
        t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.32)).text_frame
        t.clear()
        p = t.paragraphs[0]
        p.text = title
        p.font.name = FONT_MAIN
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COL[title_color]
    return sh


def bullets(shape, lines, size=20, top=0.56):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(top)
    for i, line in enumerate(lines[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(size)
        p.font.color.rgb = COL["text"]
        p.space_after = Pt(8)


# 1 cover
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 1, "Attention Is All You Need", "report-template style adapted to slides")
left = card(s, 0.7, 1.85, 8.3, 4.7, fill="bggray", title="핵심 요약")
bullets(left, [
    "Self-Attention만으로 시퀀스 변환을 수행",
    "학습 병렬화 효율을 구조적으로 크게 개선",
    "BERT/GPT 등 현대 LLM의 기반 아키텍처 제공",
], size=23, top=0.7)
right = card(s, 9.2, 1.85, 3.6, 4.7, fill="bgpurple", title="키워드")
bullets(right, ["Transformer", "Multi-Head Attention", "Positional Encoding"], size=19, top=0.72)

# 2 background
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 2, "Why this paper mattered", "from recurrence bottlenecks to attention-first modeling")
a = card(s, 0.7, 1.9, 6.1, 4.8, fill="bggray", title="Before")
b = card(s, 6.95, 1.9, 5.85, 4.8, fill="bgorange", title="After", title_color="accentorange")
bullets(a, [
    "장거리 의존성 학습이 어려움",
    "순차 처리로 병렬화 제한",
    "스케일업 시 계산 비용 급증",
], size=20)
bullets(b, [
    "Attention-only 구조로 관계 직접 계산",
    "병렬 계산 친화적 설계",
    "대규모 사전학습 모델의 토대 제공",
], size=20)

# 3 architecture
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 3, "Architecture summary", "encoder/decoder + self-attention blocks")
main = card(s, 0.7, 1.9, 12.1, 4.75, fill="bggray", title="Pipeline")
for i, lbl in enumerate(["Embedding", "Multi-Head\nAttention", "FFN", "Output"]):
    x = 1.2 + i * 2.86
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.35), Inches(2.35), Inches(1.25))
    block.fill.solid(); block.fill.fore_color.rgb = COL["dgistlight"] if i % 2 else COL["bgpurple"]
    block.line.color.rgb = COL["line"]
    tf = block.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = lbl; p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_MAIN; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["dgistnavy"]
for i in range(3):
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.58 + i * 2.86), Inches(3.73), Inches(0.45), Inches(0.38))
    arr.fill.solid(); arr.fill.fore_color.rgb = COL["dgistmid"]
    arr.line.fill.background()

# 4 equation
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 4, "Core mechanism", "scaled dot-product attention")
c = card(s, 0.7, 1.9, 12.1, 4.8, fill="bggray", title="Attention formula")
eq = s.shapes.add_textbox(Inches(1.0), Inches(2.58), Inches(11.7), Inches(0.78)).text_frame
eq.clear()
p = eq.paragraphs[0]
p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
p.font.name = FONT_MAIN; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = COL["dgistblue"]
ex = s.shapes.add_textbox(Inches(1.0), Inches(3.55), Inches(11.7), Inches(2.35)).text_frame
ex.clear()
for i, t in enumerate([
    "QK^T로 토큰 간 관련도 계산",
    "sqrt(d_k) 스케일링으로 학습 안정성 개선",
    "V 가중합으로 문맥 표현 생성",
]):
    q = ex.paragraphs[0] if i == 0 else ex.add_paragraph()
    q.text = f"• {t}"
    q.font.name = FONT_MAIN; q.font.size = Pt(22); q.font.color.rgb = COL["text"]

# 5 impact
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 5, "Results and impact", "what changed after AIAYN")
for i, (k, v, fill) in enumerate([
    ("Modeling", "Recurrence 중심에서 Attention 중심으로 전환", "bggray"),
    ("Efficiency", "병렬 학습 효율을 크게 향상", "bglight"),
    ("Legacy", "BERT/GPT 계열 LLM의 기반 확립", "bgpurple"),
]):
    x = 0.7 + i * 4.08
    box = card(s, x, 1.95, 3.9, 4.75, fill=fill, title=k)
    b = s.shapes.add_textbox(Inches(x + 0.2), Inches(2.7), Inches(3.5), Inches(3.3)).text_frame
    b.clear(); p = b.paragraphs[0]
    p.text = v
    p.font.name = FONT_MAIN; p.font.size = Pt(20); p.font.color.rgb = COL["text"]

# 6 takeaway
s = prs.slides.add_slide(prs.slide_layouts[6])
base(s, "TRANSFORMER PAPER | BRIEF")
section_title(s, 6, "Takeaways", "three points to remember")
panel = card(s, 0.7, 1.95, 12.1, 4.8, fill="bggray", title="Final summary", title_color="accentpurple")
tf = panel.text_frame
tf.clear(); tf.margin_left = Inches(0.32); tf.margin_top = Inches(0.72)
for i, txt in enumerate([
    "Transformer 이해는 LLM 이해의 출발점이다.",
    "모델 성능과 계산 효율을 함께 봐야 한다.",
    "Attention 설계 선택이 실제 성능 차이를 만든다.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{i+1}. {txt}"
    p.font.name = FONT_MAIN; p.font.size = Pt(27); p.font.bold = True; p.font.color.rgb = COL["dgistnavy"]
    p.space_after = Pt(12)

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v7.pptx"
prs.save(out)
print(out)
