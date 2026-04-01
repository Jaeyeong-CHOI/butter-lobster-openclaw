from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "최재영"
PRESENTATION_TITLE = "Attention Is All You Need 발표"
DATE_LABEL = date.today().isoformat()  # 2026-03-31

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Report-like polished palette
C_BG = RGBColor(248, 249, 252)
C_SURFACE = RGBColor(255, 255, 255)
C_LINE = RGBColor(219, 226, 243)
C_NAVY = RGBColor(14, 32, 80)
C_BLUE = RGBColor(26, 60, 140)
C_MID = RGBColor(60, 100, 180)
C_TEXT = RGBColor(34, 38, 50)
C_MUTED = RGBColor(98, 109, 138)
C_SOFT_A = RGBColor(239, 244, 255)
C_SOFT_B = RGBColor(247, 240, 255)

FONT = "Arial"


def add_base(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.26))
    top.fill.solid(); top.fill.fore_color.rgb = C_BLUE
    top.line.fill.background()

    # top-right date
    tr = slide.shapes.add_textbox(Inches(10.25), Inches(0.03), Inches(2.9), Inches(0.19)).text_frame
    tr.clear()
    p = tr.paragraphs[0]
    p.text = DATE_LABEL
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 239, 255)

    # bottom-left author
    bl = slide.shapes.add_textbox(Inches(0.82), Inches(7.05), Inches(5.8), Inches(0.2)).text_frame
    bl.clear()
    q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT
    q.font.size = Pt(11)
    q.font.color.rgb = C_MUTED

    # bottom-right title
    br = slide.shapes.add_textbox(Inches(7.2), Inches(7.05), Inches(5.2), Inches(0.2)).text_frame
    br.clear()
    r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE
    r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = C_MUTED


def add_title(slide, title, subtitle=""):
    t = slide.shapes.add_textbox(Inches(0.82), Inches(0.55), Inches(11.5), Inches(0.95)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_NAVY

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.86), Inches(1.3), Inches(11.5), Inches(0.45)).text_frame
        s.clear()
        q = s.paragraphs[0]
        q.text = subtitle
        q.font.name = FONT
        q.font.size = Pt(18)
        q.font.color.rgb = C_MUTED


def add_card(slide, x, y, w, h, fill=C_SURFACE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = C_LINE
    c.line.width = Pt(1.1)
    c.adjustments[0] = 0.05
    return c


def card_header(slide, x, y, w, text):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_BLUE


def fill_bullets(shape, bullets, size=21, top=0.58):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(0.1)
    for i, b in enumerate(bullets[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(8)


def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Attention Is All You Need", "깔끔하고 세련된 핵심 요약")

    c1 = add_card(s, 0.82, 1.95, 8.15, 4.62, C_SURFACE)
    card_header(s, 1.06, 2.12, 7.6, "Executive summary")
    fill_bullets(c1, [
        "RNN/CNN 없이 Self-Attention만으로 시퀀스 변환",
        "학습 병렬화 효율을 구조적으로 크게 개선",
        "현대 LLM(BERT/GPT) 아키텍처의 직접적 기반",
    ], size=24, top=0.68)

    c2 = add_card(s, 9.12, 1.95, 3.4, 4.62, C_SOFT_B)
    card_header(s, 9.35, 2.12, 2.9, "Keywords")
    fill_bullets(c2, ["Transformer", "Multi-Head", "Positional Encoding"], size=20, top=0.7)


def slide_why():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Why this paper mattered", "문제 정의와 구조적 전환")

    l = add_card(s, 0.82, 1.92, 6.0, 4.9, C_SURFACE)
    card_header(s, 1.06, 2.08, 5.4, "Before")
    fill_bullets(l, [
        "장거리 의존성 학습이 어려움",
        "순차 처리로 병렬화 제한",
        "확장 시 계산 비용 급증",
    ], size=21)

    r = add_card(s, 6.96, 1.92, 5.56, 4.9, C_SOFT_A)
    card_header(s, 7.2, 2.08, 4.9, "After")
    fill_bullets(r, [
        "Attention-only 구조로 관계 직접 계산",
        "병렬 계산 친화적 아키텍처",
        "LLM 시대 표준 backbone 제공",
    ], size=21)


def slide_arch():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Architecture summary", "Transformer block 구성")

    c = add_card(s, 0.82, 1.95, 12.0, 4.8, C_SURFACE)
    card_header(s, 1.06, 2.1, 11.4, "Encoder/Decoder + Attention pipeline")

    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    fills = [C_SOFT_B, C_SOFT_A, C_SOFT_B, C_SOFT_A]
    for i, lab in enumerate(labels):
        x = 1.2 + i * 2.82
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.35), Inches(2.35), Inches(1.35))
        b.fill.solid(); b.fill.fore_color.rgb = fills[i]
        b.line.color.rgb = C_LINE
        tf = b.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = lab
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY

    for i in range(3):
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.58 + i * 2.82), Inches(3.8), Inches(0.42), Inches(0.42))
        arr.fill.solid(); arr.fill.fore_color.rgb = C_MID
        arr.line.fill.background()


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Core mechanism", "Scaled Dot-Product Attention")

    c = add_card(s, 0.82, 1.95, 12.0, 4.85, C_SURFACE)
    card_header(s, 1.06, 2.1, 11.4, "Attention formula")

    eq = s.shapes.add_textbox(Inches(1.1), Inches(2.62), Inches(11.2), Inches(0.76)).text_frame
    eq.clear()
    p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT
    p.font.size = Pt(31)
    p.font.bold = True
    p.font.color.rgb = C_BLUE

    expl = s.shapes.add_textbox(Inches(1.1), Inches(3.62), Inches(11.2), Inches(2.45)).text_frame
    expl.clear()
    for i, t in enumerate([
        "QK^T로 토큰 간 관련도를 계산",
        "sqrt(d_k) 스케일링으로 학습 안정성 개선",
        "V 가중합으로 문맥 표현 생성",
    ]):
        q = expl.paragraphs[0] if i == 0 else expl.add_paragraph()
        q.text = f"• {t}"
        q.font.name = FONT
        q.font.size = Pt(22)
        q.font.color.rgb = C_TEXT
        q.space_after = Pt(10)


def slide_impact():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Results and impact", "Transformer 이후 핵심 변화")

    data = [
        ("Modeling", "Recurrence 중심에서 Attention 중심으로 전환", C_SURFACE),
        ("Efficiency", "병렬 학습 효율을 크게 향상", C_SOFT_A),
        ("Legacy", "BERT/GPT 등 현대 LLM 기반 확립", C_SOFT_B),
    ]

    for i, (k, v, fill) in enumerate(data):
        x = 0.82 + i * 4.06
        c = add_card(s, x, 2.02, 3.85, 4.68, fill)
        h = s.shapes.add_textbox(Inches(x + 0.22), Inches(2.26), Inches(3.4), Inches(0.45)).text_frame
        h.clear()
        p = h.paragraphs[0]
        p.text = k
        p.font.name = FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_BLUE

        b = s.shapes.add_textbox(Inches(x + 0.22), Inches(3.0), Inches(3.42), Inches(2.9)).text_frame
        b.clear()
        q = b.paragraphs[0]
        q.text = v
        q.font.name = FONT
        q.font.size = Pt(21)
        q.font.color.rgb = C_TEXT


def slide_end():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Takeaways", "발표 후 기억할 3가지")

    c = add_card(s, 0.82, 2.0, 12.0, 4.75, C_SURFACE)
    tf = c.text_frame
    tf.clear()
    tf.margin_left = Inches(0.33)
    tf.margin_top = Inches(0.72)

    items = [
        "Transformer 이해는 LLM 이해의 출발점이다.",
        "모델 품질과 연산 효율을 함께 봐야 한다.",
        "Attention 설계 선택이 성능 차이를 만든다.",
    ]
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {txt}"
        p.font.name = FONT
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.space_after = Pt(12)


slide_cover()
slide_why()
slide_arch()
slide_equation()
slide_impact()
slide_end()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v5.pptx"
prs.save(out)
print(out)
