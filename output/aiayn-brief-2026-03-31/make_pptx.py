from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Theme-ish constants
TITLE_COLOR = RGBColor(45, 35, 82)
BODY_COLOR = RGBColor(40, 40, 40)
ACCENT = RGBColor(76, 110, 245)

slides = [
    ("Attention Is All You Need", ["A concise brief", "Transformer의 출발점이 된 논문"]),
    ("Why this paper mattered", ["RNN/CNN 기반 Seq2Seq 한계", "장거리 의존성 + 병렬화 문제", "Attention 중심 구조로 전환"]),
    ("Core architecture", ["Encoder-Decoder Transformer", "Multi-Head Self-Attention", "FFN + Residual + LayerNorm"]),
    ("Scaled Dot-Product Attention", ["Attention(Q,K,V)=softmax(QK^T/sqrt(d_k))V", "스케일링으로 안정적 학습", "토큰 간 관계를 직접 계산"]),
    ("Multi-Head + Positional Encoding", ["여러 head가 다양한 관계를 병렬 학습", "순서 정보는 positional encoding으로 주입", "RNN 없이도 시퀀스 처리 가능"]),
    ("Results & impact", ["WMT 번역 벤치마크에서 strong performance", "학습 병렬화 효율 개선", "BERT/GPT 등 현대 LLM의 기반"]),
    ("Limitations and follow-ups", ["Self-attention O(n^2) 비용", "긴 문맥 비용 문제", "Sparse/Linear/FlashAttention 등 후속 연구"]),
    ("Takeaways", ["Transformer 이해 = LLM 이해의 핵심", "성능 + 효율을 함께 봐야 함", "Attention 설계가 실제 품질을 좌우"]),
]

for i, (title, bullets) in enumerate(slides):
    layout = prs.slide_layouts[1] if i > 0 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title
    title_tf = slide.shapes.title.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(42 if i == 0 else 36)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR

    if i == 0:
        subtitle = slide.placeholders[1]
        subtitle.text = "Summary slides for quick review"
        subtitle.text_frame.paragraphs[0].font.size = Pt(22)
        subtitle.text_frame.paragraphs[0].font.color.rgb = ACCENT
    else:
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for j, b in enumerate(bullets):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = BODY_COLOR

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief.pptx"
prs.save(out)
print(out)
