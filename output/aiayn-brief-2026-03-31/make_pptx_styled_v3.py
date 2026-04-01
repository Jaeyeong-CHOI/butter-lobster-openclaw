from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --- Canvas ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Theme tokens (aligned to report style) ---
C = {
    "bg": RGBColor(248, 249, 252),
    "surface": RGBColor(255, 255, 255),
    "line": RGBColor(216, 225, 243),
    "navy": RGBColor(14, 32, 80),
    "blue": RGBColor(26, 60, 140),
    "mid": RGBColor(60, 100, 180),
    "muted": RGBColor(94, 107, 142),
    "body": RGBColor(36, 40, 52),
    "purple_bg": RGBColor(248, 240, 255),
    "green_bg": RGBColor(235, 248, 240),
}

FONT_TITLE = "Avenir Next"
FONT_BODY = "Avenir Next"


def add_base(slide, section_label="AIAYN Brief"):
    # background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["bg"]

    # top ribbon
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.34))
    ribbon.fill.solid(); ribbon.fill.fore_color.rgb = C["blue"]
    ribbon.line.fill.background()

    # section label (top-right)
    label = slide.shapes.add_textbox(Inches(9.9), Inches(0.06), Inches(3.1), Inches(0.2)).text_frame
    label.clear()
    p = label.paragraphs[0]
    p.text = section_label
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(235, 240, 255)


def add_title(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.78), Inches(0.55), Inches(10.8), Inches(0.95)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(38)
    p.font.color.rgb = C["navy"]

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.82), Inches(1.35), Inches(11.2), Inches(0.5)).text_frame
        s.clear()
        q = s.paragraphs[0]
        q.text = subtitle
        q.font.name = FONT_BODY
        q.font.size = Pt(19)
        q.font.color.rgb = C["muted"]


def add_footer(slide, left="OpenClaw styled deck", right="Attention Is All You Need"):
    l = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(6.0), Inches(0.2)).text_frame
    l.clear()
    p = l.paragraphs[0]
    p.text = left
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.color.rgb = C["muted"]

    r = slide.shapes.add_textbox(Inches(8.3), Inches(7.08), Inches(4.2), Inches(0.2)).text_frame
    r.clear()
    q = r.paragraphs[0]
    q.text = right
    q.alignment = PP_ALIGN.RIGHT
    q.font.name = FONT_BODY
    q.font.size = Pt(11)
    q.font.color.rgb = C["muted"]


def add_card(slide, x, y, w, h, title=None, fill="surface"):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = C[fill]
    card.line.color.rgb = C["line"]
    card.line.width = Pt(1.2)
    card.adjustments[0] = 0.07  # corner radius

    if title:
        tf = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.14), Inches(w - 0.45), Inches(0.34)).text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C["blue"]
    return card


def add_bullets_to_shape(shape, bullets, size=22, top_pad=0.55):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(top_pad)
    tf.margin_bottom = Inches(0.14)

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = C["body"]
        p.space_after = Pt(10)


def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s, "Core Paper Brief")
    add_title(s, "Attention Is All You Need", "깔끔한 핵심 요약 · Transformer의 시작점")

    left = add_card(s, 0.75, 1.95, 7.1, 4.5, title="Executive summary", fill="surface")
    add_bullets_to_shape(
        left,
        [
            "RNN/CNN 대신 Self-Attention만으로 시퀀스 변환",
            "학습 병렬화 효율을 크게 개선",
            "BERT/GPT 계열 LLM 아키텍처의 기반 확립",
        ],
        size=24,
        top_pad=0.65,
    )

    right = add_card(s, 8.05, 1.95, 4.45, 4.5, title="Key terms", fill="purple_bg")
    add_bullets_to_shape(
        right,
        [
            "Transformer",
            "Multi-Head Attention",
            "Positional Encoding",
            "Scaling law 이전의 구조 혁신",
        ],
        size=20,
        top_pad=0.65,
    )

    add_footer(s)


def slide_two_col(title, subtitle, left_title, left_bullets, right_title, right_bullets, right_fill="surface"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, title, subtitle)

    l = add_card(s, 0.75, 1.95, 6.0, 4.85, title=left_title, fill="surface")
    add_bullets_to_shape(l, left_bullets, size=21)

    r = add_card(s, 6.95, 1.95, 5.55, 4.85, title=right_title, fill=right_fill)
    add_bullets_to_shape(r, right_bullets, size=21)

    add_footer(s)


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Core mechanism", "Scaled Dot-Product Attention")

    main = add_card(s, 0.75, 1.95, 12.0, 4.9, title="Attention block", fill="surface")

    # Formula line
    eq = s.shapes.add_textbox(Inches(1.1), Inches(2.75), Inches(11.2), Inches(0.8)).text_frame
    eq.clear()
    p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = "Menlo"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = C["blue"]

    note = s.shapes.add_textbox(Inches(1.1), Inches(3.65), Inches(11.2), Inches(2.5)).text_frame
    note.clear()
    points = [
        "QK^T로 토큰 간 관련도를 계산하고, softmax로 가중치를 정규화",
        "sqrt(d_k) 스케일링으로 학습 안정성과 gradient 품질 개선",
        "여러 head를 병렬 적용해 다양한 관계(문법/의미/장거리)를 동시에 포착",
    ]
    for i, t in enumerate(points):
        x = note.paragraphs[0] if i == 0 else note.add_paragraph()
        x.text = f"• {t}"
        x.font.name = FONT_BODY
        x.font.size = Pt(22)
        x.font.color.rgb = C["body"]
        x.space_after = Pt(12)

    add_footer(s)


def slide_impact_table():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(s)
    add_title(s, "Results and impact", "What changed in sequence modeling")

    card = add_card(s, 0.75, 1.95, 12.0, 4.85, title="Impact summary", fill="green_bg")

    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.24)
    tf.margin_top = Inches(0.62)

    rows = [
        ("Modeling center", "Recurrence 중심에서 Attention 중심으로 전환"),
        ("Training efficiency", "순차 의존을 줄여 병렬 학습 효율 대폭 개선"),
        ("LLM foundation", "BERT/GPT 등 현대 언어모델의 표준 backbone 확립"),
    ]

    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {k}: {v}"
        p.font.name = FONT_BODY
        p.font.size = Pt(24)
        p.font.color.rgb = C["body"]
        p.space_after = Pt(14)

    add_footer(s)


def build():
    slide_cover()
    slide_two_col(
        "Why this paper mattered",
        "문제 정의와 접근 전환",
        "Before",
        [
            "RNN Seq2Seq의 장거리 의존성 학습 한계",
            "학습이 순차적으로 진행되어 병렬화가 제한",
            "큰 데이터/모델로 확장 시 비용이 급격히 증가",
        ],
        "After",
        [
            "Attention-only 구조로 토큰 간 관계를 직접 계산",
            "병렬 계산 친화적 구조로 학습 속도/효율 개선",
            "대규모 사전학습 아키텍처의 기반 제공",
        ],
        right_fill="purple_bg",
    )

    slide_two_col(
        "Architecture summary",
        "Transformer 구성 요소",
        "Encoder / Decoder",
        [
            "블록을 반복 적층해 표현력을 확장",
            "Residual + LayerNorm으로 안정적 학습",
            "Position-wise FFN으로 비선형 변환 강화",
        ],
        "Multi-Head Self-Attention",
        [
            "서로 다른 head가 서로 다른 패턴 포착",
            "지역/전역 문맥을 동시에 반영",
            "문법적/의미적 관계를 병렬로 학습",
        ],
    )

    slide_equation()
    slide_impact_table()

    slide_two_col(
        "Limitations and follow-ups",
        "현실적 제약과 후속 발전",
        "Known limits",
        [
            "Self-Attention 비용이 O(n^2)로 증가",
            "긴 문맥에서 메모리/연산 부담이 큼",
            "실무 배포 시 latency/throughput 이슈 발생",
        ],
        "Follow-up lines",
        [
            "Sparse / Linear Attention 계열",
            "FlashAttention 기반 최적화",
            "Long-context 처리 기술 지속 발전",
        ],
        right_fill="green_bg",
    )

    slide_two_col(
        "Takeaways",
        "실무/연구 적용 관점",
        "What to remember",
        [
            "Transformer 이해는 LLM 이해의 핵심 출발점",
            "모델 품질과 계산 효율을 함께 평가해야 함",
            "Attention/Position 설계가 성능을 크게 좌우",
        ],
        "How to use",
        [
            "논문 읽기 시 구조·연산량·확장성을 함께 본다",
            "실험 설계에 효율 지표(latency/memory) 포함",
            "후속 연구와 연결해 구현 선택지 비교",
        ],
    )


build()
out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v3.pptx"
prs.save(out)
print(out)
