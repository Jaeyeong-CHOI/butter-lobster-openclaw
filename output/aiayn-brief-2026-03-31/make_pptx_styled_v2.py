from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Canvas
prs = Presentation()
prs.slide_width = int(13.333 * 914400)
prs.slide_height = int(7.5 * 914400)

# Theme tokens (aligned to LaTeX house style)
C_BG = RGBColor(247, 249, 254)
C_CARD = RGBColor(255, 255, 255)
C_BORDER = RGBColor(217, 225, 244)
C_TITLE = RGBColor(14, 32, 80)
C_ACCENT = RGBColor(26, 60, 140)
C_BODY = RGBColor(38, 41, 52)
C_SUB = RGBColor(100, 110, 138)

FONT_TITLE = "Avenir Next"
FONT_BODY = "Avenir Next"


def px(inch):
    return int(inch * 914400)


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, px(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def add_title(slide, text):
    t = slide.shapes.add_textbox(px(0.7), px(0.42), px(12.0), px(0.9)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = text
    p.font.name = FONT_TITLE
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = C_TITLE


def add_card(slide, y=1.45, h=5.55):
    card = slide.shapes.add_shape(1, px(0.7), px(y), px(12.0), px(h))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD
    card.line.color.rgb = C_BORDER
    return card


def fill_bullets(shape, bullets, size=23):
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = C_BODY


def add_footer(slide, text="OpenClaw · AIAYN brief"):
    f = slide.shapes.add_textbox(px(0.8), px(7.02), px(6.5), px(0.25)).text_frame
    f.clear()
    p = f.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.color.rgb = C_SUB


def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    # left hero card
    c = s.shapes.add_shape(1, px(0.7), px(1.25), px(7.45), px(4.9))
    c.fill.solid(); c.fill.fore_color.rgb = C_CARD
    c.line.color.rgb = C_BORDER

    tf = c.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Attention Is All You Need"
    p.font.name = FONT_TITLE
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = C_TITLE

    p2 = tf.add_paragraph()
    p2.text = "A design-aware concise briefing"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(24)
    p2.font.color.rgb = C_SUB

    p3 = tf.add_paragraph()
    p3.text = "Transformer의 출발점이 된 논문"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(24)
    p3.font.color.rgb = C_ACCENT

    # right info chip
    chip = s.shapes.add_shape(1, px(8.5), px(2.3), px(4.2), px(2.4))
    chip.fill.solid(); chip.fill.fore_color.rgb = RGBColor(237, 242, 253)
    chip.line.color.rgb = RGBColor(193, 208, 244)
    tf2 = chip.text_frame
    tf2.clear()
    q = tf2.paragraphs[0]
    q.text = "Key frame"
    q.font.size = Pt(18); q.font.bold = True; q.font.color.rgb = C_ACCENT; q.font.name = FONT_BODY
    for t in ["• Self-Attention 중심 전환", "• 병렬 학습 효율 개선", "• BERT/GPT 기반"]:
        qq = tf2.add_paragraph(); qq.text = t; qq.font.size = Pt(18); qq.font.color.rgb = C_BODY; qq.font.name = FONT_BODY

    add_footer(s)


def slide_bul(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, title)
    card = add_card(s)
    fill_bullets(card, bullets)
    add_footer(s)


def slide_equation():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Core mechanism")

    left = s.shapes.add_shape(1, px(0.7), px(1.5), px(7.3), px(5.4))
    left.fill.solid(); left.fill.fore_color.rgb = C_CARD; left.line.color.rgb = C_BORDER
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Scaled Dot-Product Attention"
    p.font.size = Pt(28); p.font.bold = True; p.font.name = FONT_TITLE; p.font.color.rgb = C_TITLE
    eq = tf.add_paragraph()
    eq.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    eq.font.size = Pt(24); eq.font.name = FONT_BODY; eq.font.color.rgb = C_ACCENT
    for t in ["• QK^T로 토큰 간 상관 계산", "• sqrt(d_k)로 안정적 softmax", "• V 가중합으로 문맥 표현 생성"]:
        b = tf.add_paragraph(); b.text = t; b.font.size = Pt(21); b.font.name = FONT_BODY; b.font.color.rgb = C_BODY

    right = s.shapes.add_shape(1, px(8.25), px(1.5), px(4.45), px(5.4))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(239, 243, 253); right.line.color.rgb = RGBColor(200, 212, 242)
    tf2 = right.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = "Why it matters"
    p2.font.size = Pt(24); p2.font.bold = True; p2.font.name = FONT_TITLE; p2.font.color.rgb = C_ACCENT
    for t in ["• recurrence 없이 관계 모델링", "• 긴 문맥 상호작용을 직접 반영", "• 현대 LLM 아키텍처의 핵심 연산"]:
        b = tf2.add_paragraph(); b.text = t; b.font.size = Pt(19); b.font.name = FONT_BODY; b.font.color.rgb = C_BODY

    add_footer(s)


def slide_table():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Results and impact")

    card = add_card(s, y=1.5, h=5.2)
    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "What changed after AIAYN"
    p.font.size = Pt(28); p.font.bold = True; p.font.name = FONT_TITLE; p.font.color.rgb = C_TITLE

    rows = [
        ("Modeling center", "Recurrence → Attention"),
        ("Training behavior", "Limited parallelism → Highly parallel"),
        ("Long-term impact", "Transformer as LLM standard backbone"),
    ]

    # pseudo-table by aligned bullets
    for k, v in rows:
        b = tf.add_paragraph()
        b.text = f"• {k}: {v}"
        b.font.size = Pt(22)
        b.font.name = FONT_BODY
        b.font.color.rgb = C_BODY

    add_footer(s)


slide_cover()
slide_bul("Why this paper mattered", [
    "RNN Seq2Seq의 장거리 의존성 학습 한계",
    "순차 처리 때문에 학습 병렬화가 어려움",
    "Attention-only 구조로 병목을 구조적으로 해결",
])
slide_bul("Architecture summary", [
    "Encoder-Decoder Transformer",
    "Multi-Head Self-Attention + Position-wise FFN",
    "Residual Connection + LayerNorm으로 안정적 학습",
])
slide_equation()
slide_table()
slide_bul("Limitations and follow-ups", [
    "Self-Attention의 O(n^2) 연산/메모리 비용",
    "긴 문맥에서 비용 급증",
    "후속: Sparse/Linear/FlashAttention, long-context 최적화",
])
slide_bul("Takeaways", [
    "Transformer 이해는 LLM 이해의 출발점",
    "성능 + 효율을 함께 최적화해야 실무 전환 가능",
    "Attention 설계 선택이 모델 품질을 크게 좌우",
])

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v2.pptx"
prs.save(out)
print(out)
