from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Clean modern tokens (report-aligned)
C_BG = RGBColor(249, 250, 253)
C_SURFACE = RGBColor(255, 255, 255)
C_LINE = RGBColor(224, 229, 242)
C_NAVY = RGBColor(14, 32, 80)
C_BLUE = RGBColor(26, 60, 140)
C_MID = RGBColor(66, 102, 178)
C_TEXT = RGBColor(34, 39, 52)
C_MUTED = RGBColor(99, 110, 141)
C_SOFT_BLUE = RGBColor(238, 244, 255)
C_SOFT_PURPLE = RGBColor(245, 239, 255)
C_SOFT_GREEN = RGBColor(237, 248, 241)

FONT = "Arial"


def add_bg(slide, tag='AIAYN Brief'):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG

    # top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_BLUE
    bar.line.fill.background()

    # tag right
    t = slide.shapes.add_textbox(Inches(10.3), Inches(0.03), Inches(2.85), Inches(0.2)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = tag
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = RGBColor(232, 238, 255)


def title_block(slide, title, subtitle=''):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(10.8), Inches(0.95)).text_frame
    tb.clear()
    p = tb.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = C_NAVY

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.83), Inches(1.3), Inches(11.5), Inches(0.45)).text_frame
        sb.clear()
        q = sb.paragraphs[0]
        q.text = subtitle
        q.font.name = FONT
        q.font.size = Pt(18)
        q.font.color.rgb = C_MUTED


def footer(slide):
    l = slide.shapes.add_textbox(Inches(0.82), Inches(7.05), Inches(6), Inches(0.2)).text_frame
    l.clear()
    p = l.paragraphs[0]
    p.text = 'OpenClaw styled deck'
    p.font.name = FONT; p.font.size = Pt(11); p.font.color.rgb = C_MUTED

    r = slide.shapes.add_textbox(Inches(8.0), Inches(7.05), Inches(4.5), Inches(0.2)).text_frame
    r.clear()
    q = r.paragraphs[0]
    q.text = 'Attention Is All You Need'
    q.alignment = PP_ALIGN.RIGHT
    q.font.name = FONT; q.font.size = Pt(11); q.font.color.rgb = C_MUTED


def card(slide, x, y, w, h, fill=C_SURFACE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = C_LINE
    c.line.width = Pt(1.0)
    c.adjustments[0] = 0.05
    return c


def card_title(slide, x, y, w, text):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35)).text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(15.5); p.font.bold = True; p.font.color.rgb = C_BLUE


def write_bullets(shape, items, size=21, top=0.54):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(0.12)
    for i, item in enumerate(items[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {item}'
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(8)


# 1) Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'Core Paper Brief')
title_block(s, 'Attention Is All You Need', '깔끔한 핵심 요약 · Transformer의 시작점')

left = card(s, 0.8, 1.95, 8.2, 4.65, C_SURFACE)
card_title(s, 1.05, 2.1, 7.6, 'Executive summary')
write_bullets(left, [
    'RNN/CNN 대신 Self-Attention으로 시퀀스 변환',
    '학습 병렬화 효율을 구조적으로 개선',
    'BERT/GPT 계열의 기반 아키텍처 확립'
], size=24, top=0.68)

right = card(s, 9.15, 1.95, 3.35, 4.65, C_SOFT_PURPLE)
card_title(s, 9.4, 2.1, 2.8, 'Key terms')
write_bullets(right, ['Transformer', 'Multi-Head Attention', 'Positional Encoding'], size=20, top=0.7)
footer(s)

# 2) Why mattered (before/after)
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_block(s, 'Why this paper mattered', '문제 정의와 구조적 전환')

b = card(s, 0.8, 1.9, 6.0, 4.9, C_SURFACE)
card_title(s, 1.05, 2.05, 5.5, 'Before')
write_bullets(b, [
    'RNN Seq2Seq는 장거리 의존성 학습이 어려움',
    '순차 처리로 학습 병렬화가 제한됨',
    '대규모 확장에서 시간/비용 병목 발생'
], size=20)

a = card(s, 6.95, 1.9, 5.55, 4.9, C_SOFT_BLUE)
card_title(s, 7.2, 2.05, 5.0, 'After')
write_bullets(a, [
    'Attention-only 구조로 관계를 직접 계산',
    '병렬 계산 친화적 구조로 효율 개선',
    '현대 LLM 아키텍처의 표준 기반 제공'
], size=20)
footer(s)

# 3) Architecture diagram style
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_block(s, 'Architecture summary', 'Transformer block 구성')

main = card(s, 0.8, 1.95, 12.0, 4.8, C_SURFACE)

# blocks
for i, label in enumerate(['Input Embedding', 'Multi-Head Attention', 'Feed-Forward', 'Output']):
    x = 1.2 + i * 2.8
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.35), Inches(2.35), Inches(1.35))
    block.fill.solid(); block.fill.fore_color.rgb = C_SOFT_BLUE if i % 2 else C_SOFT_PURPLE
    block.line.color.rgb = C_LINE
    tf = block.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = C_NAVY

# arrows
for i in range(3):
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.55 + i*2.8), Inches(3.78), Inches(0.45), Inches(0.45))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_MID
    arr.line.fill.background()

note = s.shapes.add_textbox(Inches(1.1), Inches(2.35), Inches(11.3), Inches(0.6)).text_frame
note.clear()
p = note.paragraphs[0]
p.text = 'Encoder/Decoder는 Self-Attention + FFN + Residual/LayerNorm 블록을 반복해 표현력을 확장한다.'
p.font.name = FONT; p.font.size = Pt(20); p.font.color.rgb = C_TEXT
footer(s)

# 4) Equation slide
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_block(s, 'Core mechanism', 'Scaled Dot-Product Attention')

c = card(s, 0.8, 1.95, 12.0, 4.85, C_SURFACE)
eq = s.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(11.3), Inches(0.8)).text_frame
eq.clear()
p = eq.paragraphs[0]
p.text = 'Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V'
p.font.name = FONT; p.font.size = Pt(31); p.font.bold = True; p.font.color.rgb = C_BLUE

write = s.shapes.add_textbox(Inches(1.1), Inches(3.6), Inches(11.2), Inches(2.6)).text_frame
write.clear()
for i, t in enumerate([
    'QK^T로 토큰 간 관련도 계산',
    'sqrt(d_k)로 softmax 포화 완화 및 안정화',
    'V 가중합으로 문맥 임베딩 생성'
]):
    pp = write.paragraphs[0] if i == 0 else write.add_paragraph()
    pp.text = f'• {t}'
    pp.font.name = FONT; pp.font.size = Pt(22); pp.font.color.rgb = C_TEXT
footer(s)

# 5) Impact cards
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_block(s, 'Results and impact', 'Transformer 이후 변화')

for i, (t, v) in enumerate([
    ('Modeling', 'Recurrence → Attention 중심 전환'),
    ('Efficiency', '병렬 학습 효율 대폭 개선'),
    ('Legacy', 'BERT/GPT 등 현대 LLM의 기반')
]):
    x = 0.8 + i*4.17
    cc = card(s, x, 2.05, 3.95, 4.55, C_SOFT_GREEN if i==1 else C_SURFACE)
    tt = s.shapes.add_textbox(Inches(x+0.25), Inches(2.3), Inches(3.4), Inches(0.5)).text_frame
    tt.clear(); p = tt.paragraphs[0]
    p.text = t
    p.font.name = FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = C_BLUE
    bb = s.shapes.add_textbox(Inches(x+0.25), Inches(3.0), Inches(3.45), Inches(2.9)).text_frame
    bb.clear(); q = bb.paragraphs[0]
    q.text = v
    q.font.name = FONT; q.font.size = Pt(21); q.font.color.rgb = C_TEXT
footer(s)

# 6) Limitations
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
title_block(s, 'Limitations and follow-ups', '현실적 제약과 개선 방향')

l = card(s, 0.8, 2.0, 6.0, 4.8, C_SURFACE)
card_title(s, 1.05, 2.15, 5.4, 'Known limits')
write_bullets(l, [
    'Self-Attention 계산량 O(n^2)',
    '긴 문맥에서 메모리/연산 부담 증가',
    '실무 latency/throughput 최적화 필요'
], size=21)

r = card(s, 6.95, 2.0, 5.55, 4.8, C_SOFT_BLUE)
card_title(s, 7.2, 2.15, 4.9, 'Follow-up lines')
write_bullets(r, [
    'Sparse / Linear Attention',
    'FlashAttention 최적화',
    'Long-context 특화 아키텍처'
], size=21)
footer(s)

# 7) Takeaways
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'Final Takeaways')
title_block(s, 'Takeaways', '발표 후 기억해야 할 3가지')

panel = card(s, 0.8, 2.0, 12.0, 4.75, C_SURFACE)
tf = panel.text_frame
tf.clear(); tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.7)
for i, txt in enumerate([
    'Transformer 이해는 LLM 이해의 출발점이다.',
    '성능과 효율(연산/메모리)을 함께 봐야 한다.',
    'Attention 설계 선택이 실제 품질을 좌우한다.'
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'{i+1}. {txt}'
    p.font.name = FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_NAVY
    p.space_after = Pt(14)

footer(s)

out = '/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_styled_v4.pptx'
prs.save(out)
print(out)
