from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

AUTHOR = "Jaeyeong CHOI"
PRESENTATION_TITLE = "Attention Is All You Need Presentation"
DATE_LABEL = date.today().isoformat()

# A4 portrait reference (mm)
A4_W = 210.0
A4_H = 297.0

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width / 914400.0
SH = prs.slide_height / 914400.0

COL = {
    "navy": RGBColor(14, 32, 80),
    "blue": RGBColor(26, 60, 140),
    "mid": RGBColor(60, 100, 180),
    "light": RGBColor(220, 230, 250),
    "bg": RGBColor(248, 249, 252),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(216, 225, 243),
    "text": RGBColor(36, 40, 52),
    "muted": RGBColor(98, 109, 138),
    "soft1": RGBColor(248, 240, 255),
    "soft2": RGBColor(255, 247, 235),
}
FONT = "Apple SD Gothic Neo"


def xm(mm):
    return Inches((mm / A4_W) * SW)


def ym(mm):
    return Inches((mm / A4_H) * SH)


def add_meta(slide, dark=False):
    tr = slide.shapes.add_textbox(xm(157), ym(2), xm(50), ym(7)).text_frame
    tr.clear()
    p = tr.paragraphs[0]
    p.text = DATE_LABEL
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 238, 255) if dark else COL["muted"]

    bl = slide.shapes.add_textbox(xm(7), ym(286), xm(85), ym(8)).text_frame
    bl.clear()
    q = bl.paragraphs[0]
    q.text = AUTHOR
    q.font.name = FONT
    q.font.size = Pt(11)
    q.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]

    br = slide.shapes.add_textbox(xm(95), ym(286), xm(108), ym(8)).text_frame
    br.clear()
    r = br.paragraphs[0]
    r.text = PRESENTATION_TITLE
    r.alignment = PP_ALIGN.RIGHT
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(220, 230, 250) if dark else COL["muted"]


def cover_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COL["navy"]

    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xm(0), ym(0), xm(210), ym(11))
    top.fill.solid(); top.fill.fore_color.rgb = COL["blue"]; top.line.fill.background()

    hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xm(15), ym(48), xm(180), ym(182))
    hero.fill.solid(); hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.fill.transparency = 0.90
    hero.line.color.rgb = RGBColor(165, 188, 238)

    tf = hero.text_frame
    tf.clear(); tf.margin_left = xm(8); tf.margin_top = ym(12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT; p.font.size = Pt(48); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    q = tf.add_paragraph()
    q.text = subtitle
    q.font.name = FONT; q.font.size = Pt(20)
    q.font.color.rgb = COL["light"]

    k = tf.add_paragraph()
    k.text = "A4-style design remapped to 16:9 slide canvas"
    k.font.name = FONT; k.font.size = Pt(15)
    k.font.color.rgb = RGBColor(195, 212, 250)

    add_meta(s, dark=True)


def body_header(slide, idx, title, subtitle=""):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = COL["bg"]

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xm(0), ym(0), xm(210), ym(10))
    bar.fill.solid(); bar.fill.fore_color.rgb = COL["blue"]; bar.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xm(10), ym(21), xm(10), ym(16))
    marker.fill.solid(); marker.fill.fore_color.rgb = COL["blue"]; marker.line.fill.background()
    mt = marker.text_frame; mt.clear(); mp = mt.paragraphs[0]
    mp.text = str(idx)
    mp.alignment = PP_ALIGN.CENTER
    mp.font.name = FONT; mp.font.size = Pt(14); mp.font.bold = True
    mp.font.color.rgb = RGBColor(255, 255, 255)

    tt = slide.shapes.add_textbox(xm(23), ym(20), xm(170), ym(20)).text_frame
    tt.clear(); t = tt.paragraphs[0]
    t.text = title
    t.font.name = FONT; t.font.size = Pt(31); t.font.bold = True
    t.font.color.rgb = COL["navy"]

    u = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xm(23), ym(39), xm(170), ym(1.8))
    u.fill.solid(); u.fill.fore_color.rgb = COL["blue"]; u.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(xm(24), ym(42), xm(167), ym(11)).text_frame
        st.clear(); sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT; sp.font.size = Pt(15)
        sp.font.color.rgb = COL["muted"]

    add_meta(slide, dark=False)


def card(slide, x, y, w, h, title=None, fill=RGBColor(255,255,255)):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xm(x), ym(y), xm(w), ym(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = COL["line"]; c.line.width = Pt(1.0)
    c.adjustments[0] = 0.05
    if title:
        th = slide.shapes.add_textbox(xm(x+4), ym(y+4), xm(w-8), ym(10)).text_frame
        th.clear(); p = th.paragraphs[0]
        p.text = title
        p.font.name = FONT; p.font.size = Pt(14.5); p.font.bold = True
        p.font.color.rgb = COL["blue"]
    return c


def fill_bullets(shape, items, size=19):
    tf = shape.text_frame
    tf.clear(); tf.margin_left = xm(4); tf.margin_right = xm(3); tf.margin_top = ym(18)
    for i, it in enumerate(items[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {it}"
        p.font.name = FONT; p.font.size = Pt(size); p.font.color.rgb = COL["text"]
        p.space_after = Pt(8)


def slide_context():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 1, "Problem and context", "from recurrence bottlenecks to attention-first modeling")
    l = card(s, 10, 58, 96, 188, "Before", COL["card"])
    r = card(s, 108, 58, 92, 188, "After", COL["soft2"])
    fill_bullets(l, [
        "Long-range dependency learning was difficult",
        "Sequential processing limited parallelism",
        "Scaling increased training cost significantly",
    ])
    fill_bullets(r, [
        "Self-attention directly models token relations",
        "Parallel-friendly architecture improved throughput",
        "Set foundation for large-scale language models",
    ])


def slide_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 2, "Architecture overview", "encoder/decoder with repeated attention blocks")
    panel = card(s, 10, 58, 190, 188, "Transformer pipeline", COL["card"])

    labels = ["Embedding", "Multi-Head\nAttention", "FFN", "Output"]
    fills = [COL["soft1"], RGBColor(236,248,241), COL["soft1"], RGBColor(236,248,241)]
    for i, lbl in enumerate(labels):
        x = 18 + i * 44
        b = slide = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xm(x), ym(114), xm(34), ym(48))
        b.fill.solid(); b.fill.fore_color.rgb = fills[i]
        b.line.color.rgb = COL["line"]
        tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = lbl; p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COL["navy"]

    for i in range(3):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, xm(53 + i*44), ym(128), xm(8), ym(14))
        a.fill.solid(); a.fill.fore_color.rgb = COL["mid"]
        a.line.fill.background()


def slide_mechanism():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 3, "Core mechanism", "scaled dot-product attention")
    panel = card(s, 10, 58, 190, 188, "Attention formula", COL["card"])

    eq = s.shapes.add_textbox(xm(16), ym(90), xm(178), ym(24)).text_frame
    eq.clear(); p = eq.paragraphs[0]
    p.text = "Attention(Q, K, V) = softmax(QK^T / sqrt(d_k))V"
    p.font.name = FONT; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = COL["blue"]

    ex = s.shapes.add_textbox(xm(16), ym(124), xm(178), ym(104)).text_frame
    ex.clear()
    for i, t in enumerate([
        "QK^T computes token-to-token relevance",
        "sqrt(d_k) scaling stabilizes optimization",
        "Weighted V aggregation builds contextual representations",
    ]):
        q = ex.paragraphs[0] if i == 0 else ex.add_paragraph()
        q.text = f"• {t}"
        q.font.name = FONT; q.font.size = Pt(20); q.font.color.rgb = COL["text"]


def slide_results():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 4, "Results and impact", "A4-style table block remapped to slide")
    panel = card(s, 10, 58, 190, 188, None, COL["card"])

    table = s.shapes.add_table(4, 2, xm(16), ym(82), xm(178), ym(154)).table
    table.columns[0].width = xm(45)
    table.columns[1].width = xm(133)

    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Impact"
    for c in range(2):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = COL["blue"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)

    rows = [
        ("Modeling", "Shift from recurrence-centric to attention-centric sequence modeling"),
        ("Efficiency", "Substantially improved training parallelism and scalability"),
        ("Legacy", "Established backbone used by BERT, GPT, and modern LLM families"),
    ]
    for i, (k, v) in enumerate(rows, start=1):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
        for c in range(2):
            p = table.cell(i, c).text_frame.paragraphs[0]
            p.font.name = FONT; p.font.size = Pt(13.5); p.font.color.rgb = COL["text"]


def slide_takeaways():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_header(s, 5, "Limitations and takeaways", "what to remember for modern LLM design")
    l = card(s, 10, 58, 96, 188, "Limitations", COL["card"])
    r = card(s, 108, 58, 92, 188, "Takeaways", COL["soft1"])
    fill_bullets(l, [
        "Self-attention cost still grows quadratically",
        "Long-context inference increases memory and latency",
        "Production deployment needs careful efficiency tuning",
    ], 19)
    fill_bullets(r, [
        "Transformer literacy is core LLM literacy",
        "Quality and efficiency should be optimized together",
        "Attention design choices strongly shape capability",
    ], 19)


cover_slide("Attention Is All You Need", "A4-to-16:9 remapped design")
slide_context()
slide_architecture()
slide_mechanism()
slide_results()
slide_takeaways()

out = "/Users/jaeyeong_openclaw/.openclaw/workspace/output/aiayn-brief-2026-03-31/attention_is_all_you_need_brief_a4-remap_v13.pptx"
prs.save(out)
print(out)
